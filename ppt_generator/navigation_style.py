"""Apply project-level navigation styling to compatible PPT layouts."""

from __future__ import annotations

from collections.abc import Sequence
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from .project import PresentationStyle, default_navigation_items


TEMPLATE_NAVIGATION_ITEMS = (
    "企业简介",
    "工艺分析",
    "设备布局",
    "模块介绍",
    "供货范围",
)
NAVIGATION_ITEMS = tuple(item.name for item in default_navigation_items())
NAVIGATION_FONT = "Microsoft YaHei"
NAVIGATION_LEFT = 0.0
NAVIGATION_RIGHT = 11.55
SLIDE_WIDTH = 13.333333


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.removeprefix("#"))


def _is_dark(value: str) -> bool:
    channels = [int(value[index : index + 2], 16) / 255 for index in (1, 3, 5)]
    linear = [
        channel / 12.92
        if channel <= 0.04045
        else ((channel + 0.055) / 1.055) ** 2.4
        for channel in channels
    ]
    luminance = 0.2126 * linear[0] + 0.7152 * linear[1] + 0.0722 * linear[2]
    return luminance < 0.32


def _set_typeface(run: Any) -> None:
    run.font.name = NAVIGATION_FONT
    properties = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = properties.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            properties.append(node)
        node.set("typeface", NAVIGATION_FONT)


def _style_navigation_text(
    shape: Any,
    color: RGBColor,
    font_size: float,
) -> None:
    frame = shape.text_frame
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_before = 0
        paragraph.space_after = 0
        for run in paragraph.runs:
            _set_typeface(run)
            run.font.size = Pt(font_size)
            run.font.bold = True
            run.font.color.rgb = color


def _layout_navigation_shapes(layout: Any) -> tuple[Any, Any, list[Any]] | None:
    navigation_shapes = [shape for shape in layout.shapes if shape.top < Inches(0.95)]
    text_shapes = sorted(
        (
            shape
            for shape in navigation_shapes
            if getattr(shape, "has_text_frame", False)
            and shape.text.strip() in TEMPLATE_NAVIGATION_ITEMS
        ),
        key=lambda shape: shape.left,
    )
    if len(text_shapes) != len(TEMPLATE_NAVIGATION_ITEMS):
        return None
    background = next(
        (
            shape
            for shape in navigation_shapes
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and shape.width > Inches(12.0)
        ),
        None,
    )
    active_bar = next(
        (
            shape
            for shape in navigation_shapes
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and shape is not background
        ),
        None,
    )
    if background is None or active_bar is None:
        return None
    return background, active_bar, text_shapes


def _style_layout(layout: Any, style: PresentationStyle) -> int | None:
    found = _layout_navigation_shapes(layout)
    if found is None:
        return None
    background, active_bar, text_shapes = found
    navigation_shapes = [shape for shape in layout.shapes if shape.top < Inches(0.95)]
    active_center = active_bar.left + active_bar.width // 2
    default_active_index = min(
        range(len(text_shapes)),
        key=lambda index: abs(
            text_shapes[index].left + text_shapes[index].width // 2 - active_center
        ),
    )
    background.top = 0
    background.height = Inches(style.navigation_height)
    background.fill.solid()
    background.fill.fore_color.rgb = _rgb("#FFFFFF")
    background.line.fill.background()
    effect_list = background.element.spPr.find(qn("a:effectLst"))
    if effect_list is not None:
        for outer_shadow in list(effect_list.findall(qn("a:outerShdw"))):
            effect_list.remove(outer_shadow)
        if len(effect_list) == 0:
            background.element.spPr.remove(effect_list)

    active_bar.fill.solid()
    active_bar.fill.fore_color.rgb = _rgb("#FFFFFF")
    active_bar.line.fill.background()

    for shape in text_shapes:
        shape.text = ""

    for shape in navigation_shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            shape.line.fill.background()
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            aspect = shape.width / shape.height
            target_height_inches = style.navigation_height * 0.50
            target_height = Inches(target_height_inches)
            target_width = int(target_height * aspect)
            shape.height = target_height
            shape.width = target_width
            shape.left = Inches(13.06) - target_width
            shape.top = Inches(
                (style.navigation_height - target_height_inches) / 2
            )
    return default_active_index


def _add_slide_navigation(
    slide: Any,
    style: PresentationStyle,
    active_index: int | None,
) -> None:
    item_count = len(style.navigation_items)
    cell_width = (NAVIGATION_RIGHT - NAVIGATION_LEFT) / item_count
    dark_active_background = _is_dark(style.navigation_background)
    inactive_color = _rgb("#515960")
    active_color = _rgb("#FFFFFF" if dark_active_background else "#C90000")
    separator_color = _rgb("#D3D9DE")
    baseline_color = _rgb("#D3D9DE")
    active_strip_color = _rgb("#C90000")
    font_size = style.resolved_navigation_font_size()

    if active_index is not None and 0 <= active_index < item_count:
        active_background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(NAVIGATION_LEFT + cell_width * active_index),
            0,
            Inches(cell_width),
            Inches(style.navigation_height),
        )
        active_background.name = "KY_NAV_ACTIVE_BACKGROUND"
        active_background.fill.solid()
        active_background.fill.fore_color.rgb = _rgb(style.navigation_background)
        active_background.line.fill.background()

    for index in range(1, item_count):
        boundary = NAVIGATION_LEFT + cell_width * index
        separator = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(boundary),
            Inches(style.navigation_height * 0.28),
            Pt(0.6),
            Inches(style.navigation_height * 0.44),
        )
        separator.name = f"KY_NAV_SEPARATOR_{index}"
        separator.fill.solid()
        separator.fill.fore_color.rgb = separator_color
        separator.line.fill.background()

    baseline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(NAVIGATION_LEFT),
        Inches(style.navigation_height - 0.04),
        Inches(SLIDE_WIDTH - NAVIGATION_LEFT),
        Inches(0.04),
    )
    baseline.name = "KY_NAV_BASELINE"
    baseline.fill.solid()
    baseline.fill.fore_color.rgb = baseline_color
    baseline.line.fill.background()

    if active_index is not None and 0 <= active_index < item_count:
        active_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(NAVIGATION_LEFT + cell_width * active_index),
            Inches(style.navigation_height - 0.04),
            Inches(cell_width),
            Inches(0.04),
        )
        active_bar.name = "KY_NAV_ACTIVE"
        active_bar.fill.solid()
        active_bar.fill.fore_color.rgb = active_strip_color
        active_bar.line.fill.background()

    for index, item in enumerate(style.navigation_items):
        text_shape = slide.shapes.add_textbox(
            Inches(NAVIGATION_LEFT + cell_width * index),
            0,
            Inches(cell_width),
            Inches(style.navigation_height - 0.03),
        )
        text_shape.name = f"KY_NAV_TEXT_{index + 1}"
        text_shape.text = item.name.strip()
        _style_navigation_text(
            text_shape,
            active_color if index == active_index else inactive_color,
            font_size,
        )


def apply_navigation_style(
    presentation: Any,
    style: PresentationStyle,
    active_indices: Sequence[int | None] | None = None,
) -> int:
    """Apply configured navigation to compatible layouts and generated slides."""
    style.validate()
    slides = list(presentation.slides)
    if active_indices is not None and len(active_indices) != len(slides):
        raise ValueError("导航活动栏目数量必须与幻灯片数量一致")

    layout_defaults: dict[str, int] = {}
    for layout in presentation.slide_layouts:
        default_index = _style_layout(layout, style)
        if default_index is not None:
            layout_defaults[str(layout.part.partname)] = default_index

    for slide_index, slide in enumerate(slides):
        layout_key = str(slide.slide_layout.part.partname)
        if layout_key not in layout_defaults:
            continue
        active_index = (
            active_indices[slide_index]
            if active_indices is not None
            else min(layout_defaults[layout_key], len(style.navigation_items) - 1)
        )
        _add_slide_navigation(slide, style, active_index)
    return len(layout_defaults)
