"""Parse a visual optical FAR workbook and map it into PPT project modules."""

from __future__ import annotations

import hashlib
import re
from collections import OrderedDict
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from openpyxl import load_workbook
from pptx import Presentation

from .project import AssetRecord, PptProject, ProjectModule, ProjectSlide, SourceRecord
from .template_renderer import TemplateManifest


class OpticalFarError(ValueError):
    """Raised when an optical FAR cannot be parsed or applied safely."""


@dataclass(frozen=True)
class FarRequirement:
    category: str
    target: str
    standard: str
    mark: str = ""


@dataclass(frozen=True)
class FarImage:
    data: bytes
    extension: str
    caption: str = ""


@dataclass
class FarStation:
    station: str
    view: str
    config: str
    items: list[str]
    camera: str = ""
    lens: str = ""
    lights: list[str] = field(default_factory=list)
    images: list[FarImage] = field(default_factory=list)


@dataclass
class OpticalFarData:
    source_path: Path
    project_code: str
    production_rate: str
    special_notes: str
    requirements: list[FarRequirement]
    stations: list[FarStation]

    @property
    def image_count(self) -> int:
        return sum(len(station.images) for station in self.stations)


@dataclass(frozen=True)
class OpticalFarApplyResult:
    effect_pages: int
    requirements: int
    stations: int
    asset_directory: Path
    camera_summary: str
    lens_summary: str
    light_summary: str


_EMPTY_VALUES = {"", "/", "／", "无", "None", "none"}
_COUNT_PATTERN = re.compile(r"(?:\*|x|X|×)\s*(\d+)")


def _text(value: Any) -> str:
    if value is None:
        return ""
    return str(value).strip()


def _clean_label(value: Any) -> str:
    return _text(value).rstrip(":：;；").strip()


def _cell_text(worksheet: Any, row: int, column: int) -> str:
    return _text(worksheet.cell(row, column).value)


def _value_after_label(worksheet: Any, labels: set[str]) -> str:
    for row in range(1, int(worksheet.max_row or 0) + 1):
        for column in range(1, int(worksheet.max_column or 0) + 1):
            if _clean_label(worksheet.cell(row, column).value) not in labels:
                continue
            for candidate in range(column + 1, min(column + 14, worksheet.max_column) + 1):
                value = _cell_text(worksheet, row, candidate)
                if value:
                    return "" if value in _EMPTY_VALUES else value
    return ""


def _normalize_mark(value: str) -> str:
    mark = _text(value)
    return {"？": "?", "×": "X", "x": "X"}.get(mark, mark)


def _parse_requirements(worksheet: Any) -> list[FarRequirement]:
    max_row = int(worksheet.max_row or 0)
    max_column = int(worksheet.max_column or 0)
    title_row = 0
    title_column = 0
    for row in range(1, max_row + 1):
        for column in range(1, max_column + 1):
            if _cell_text(worksheet, row, column) == "检测项汇总":
                title_row, title_column = row, column
                break
        if title_row:
            break
    if not title_row:
        raise OpticalFarError("未找到右上角“检测项汇总”")

    standard_column = next(
        (
            column
            for column in range(title_column + 1, max_column + 1)
            if "验收标准" in _cell_text(worksheet, title_row, column)
        ),
        0,
    )
    if not standard_column:
        raise OpticalFarError("检测项汇总缺少“验收标准”列")

    header_row = next(
        (
            row
            for row in range(title_row + 1, min(title_row + 8, max_row) + 1)
            if "检测项" in _cell_text(worksheet, row, title_column)
        ),
        title_row + 2,
    )
    category_column = title_column
    target_column = title_column + 2
    material_columns = [
        column
        for column in range(target_column + 1, standard_column)
        if _cell_text(worksheet, title_row, column)
    ]
    if not material_columns:
        material_columns = list(range(target_column + 1, standard_column))

    requirements: list[FarRequirement] = []
    current_category = ""
    for row in range(header_row + 1, max_row + 1):
        row_text = " ".join(
            _cell_text(worksheet, row, column)
            for column in range(1, max_column + 1)
            if _cell_text(worksheet, row, column)
        )
        if requirements and (
            "√表示传统算法进行检测" in row_text
            or "3.检测项分析" in row_text
            or ("工位视角" in row_text and "检测项目" in row_text)
        ):
            break
        category_cell = _cell_text(worksheet, row, category_column)
        if category_cell:
            current_category = category_cell
        target = _cell_text(worksheet, row, target_column)
        standard = _cell_text(worksheet, row, standard_column)
        if not standard:
            for column in range(standard_column + 1, min(max_column, standard_column + 10) + 1):
                standard = _cell_text(worksheet, row, column)
                if standard:
                    break
        marks = [
            _normalize_mark(_cell_text(worksheet, row, column))
            for column in material_columns
        ]
        mark = next((value for value in marks if value), "")
        if not target and not standard and not mark:
            continue
        if not current_category and not target:
            continue
        if not standard and not mark:
            continue
        requirements.append(
            FarRequirement(
                category=current_category,
                target=target,
                standard=standard,
                mark=mark,
            )
        )
    if not requirements:
        raise OpticalFarError("检测项汇总中没有可用记录")
    return requirements


def _find_analysis_header(worksheet: Any) -> tuple[int, dict[str, int]]:
    for row in range(1, int(worksheet.max_row or 0) + 1):
        columns: dict[str, int] = {}
        for column in range(1, int(worksheet.max_column or 0) + 1):
            value = _cell_text(worksheet, row, column)
            if "工位视角" in value:
                columns["view"] = column
            elif "工位配置" in value:
                columns["config"] = column
            elif "检测项目及精度" in value or value == "检测项目":
                columns["items"] = column
            elif "软件截图" in value or "成像效果" in value:
                columns["images"] = column
        if {"view", "config", "items", "images"} <= set(columns):
            return row, columns
    raise OpticalFarError("未找到“工位视角/工位配置/检测项目/软件截图”表头")


def _config_parts(config: str) -> tuple[str, str, list[str]]:
    camera = ""
    lens = ""
    lights: list[str] = []
    in_light_block = False
    for raw_line in str(config or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        if re.match(r"^相机\s*[:：]", line):
            camera = re.split(r"[:：]", line, maxsplit=1)[1].strip()
            in_light_block = False
            continue
        if re.match(r"^镜头\s*[:：]", line):
            lens = re.split(r"[:：]", line, maxsplit=1)[1].strip()
            in_light_block = False
            continue
        if re.match(r"^光源\s*[:：]", line):
            value = re.split(r"[:：]", line, maxsplit=1)[1].strip()
            if value and not re.fullmatch(r"(?:\*|x|X|×)?\s*\d*", value):
                lights.append(value)
            in_light_block = True
            continue
        if in_light_block and ("光" in line or "灯" in line):
            lights.append(line)
    return camera, lens, list(dict.fromkeys(lights))


def _image_extension(image: Any, data: bytes) -> str:
    extension = str(getattr(image, "format", "") or "").lower().lstrip(".")
    if extension in {"jpeg", "jpg", "png", "bmp", "gif", "tiff", "webp"}:
        return "jpg" if extension == "jpeg" else extension
    if data.startswith(b"\x89PNG\r\n\x1a\n"):
        return "png"
    if data.startswith(b"\xff\xd8"):
        return "jpg"
    return "bin"


def _valid_caption(value: str) -> bool:
    text = str(value or "").strip()
    return bool(
        text
        and not text.startswith("=")
        and "DISPIMG" not in text
        and text not in {"效果图", "软件截图", "产品图片"}
    )


def _image_caption(
    worksheet: Any,
    image_row: int,
    image_column: int,
    block_end: int,
) -> str:
    best: tuple[tuple[int, int], str] | None = None
    for row in range(image_row + 1, block_end):
        for offset in (0, 1, -1, 2, -2):
            column = image_column + offset
            if column < 1:
                continue
            value = _cell_text(worksheet, row, column)
            if not _valid_caption(value):
                continue
            score = (abs(offset), row - image_row)
            if best is None or score < best[0]:
                best = (score, value)
    return best[1] if best else ""


def _normalized_item_text(value: str) -> str:
    return re.sub(r"\s+", "", str(value or "").strip()).translate(
        str.maketrans({"（": "(", "）": ")", "，": ",", "：": ":"})
    )


def _unique_station_items(station: FarStation) -> list[str]:
    items: list[str] = []
    seen: set[str] = set()
    for item in station.items:
        text = str(item or "").strip()
        key = _normalized_item_text(text)
        if not key or key in seen:
            continue
        seen.add(key)
        items.append(text)
    return items


def _matching_station_image_index(
    station: FarStation,
    item: str,
) -> int | None:
    item_key = _normalized_item_text(item)
    for index, image in enumerate(station.images):
        caption_key = _normalized_item_text(image.caption)
        if not caption_key or caption_key.upper() == "OK":
            continue
        if caption_key == item_key or caption_key in item_key or item_key in caption_key:
            return index
    return None


def _ok_station_image_index(station: FarStation) -> int | None:
    for index, image in enumerate(station.images):
        if _normalized_item_text(image.caption).upper() == "OK":
            return index
    return None


def _concise_item_name(item: str) -> str:
    return re.split(r"[（(]", str(item or "").strip(), maxsplit=1)[0].strip()


def _effect_item_text(items: list[str]) -> str:
    names = [_concise_item_name(item) or item for item in items]
    lines: list[str] = []
    for name in names:
        prefix = "检测项：" if not lines else ""
        if not lines:
            lines.append(prefix + name)
            continue
        candidate = lines[-1] + "、" + name
        if len(candidate) <= 14:
            lines[-1] = candidate
        elif len(lines) == 1:
            lines.append(name)
        else:
            raise OpticalFarError("OK页检测项超过两行容量，请拆分工位或调整模板")
    return "\n".join(lines)


def _ok_effect_item_text(items: list[str]) -> str:
    names = [_concise_item_name(item) or item for item in items]
    item_list = "、".join(names)
    if len(item_list) + 1 > 14:
        raise OpticalFarError("OK页括号内检测项超过单行容量，请调整模板")
    return f"OK样件图（\n{item_list}）"


def _parse_stations(worksheet: Any) -> list[FarStation]:
    header_row, columns = _find_analysis_header(worksheet)
    group_rows: list[tuple[int, str, str, list[str]]] = []
    for row in range(header_row + 1, int(worksheet.max_row or 0) + 1):
        view = _cell_text(worksheet, row, columns["view"])
        config = _cell_text(worksheet, row, columns["config"])
        item_text = _cell_text(worksheet, row, columns["items"])
        items = [item.strip() for item in item_text.splitlines() if item.strip()]
        if view and config and items:
            group_rows.append((row, view, config, items))
    if not group_rows:
        raise OpticalFarError("检测项分析中没有可用工位")

    anchored_images: list[tuple[int, int, Any]] = []
    for image in list(getattr(worksheet, "_images", []) or []):
        anchor = getattr(image, "anchor", None)
        marker = getattr(anchor, "_from", None)
        if marker is None:
            continue
        anchored_images.append((int(marker.row) + 1, int(marker.col) + 1, image))

    stations: list[FarStation] = []
    image_floor = columns["images"]
    final_end = max(int(worksheet.max_row or 0), max((item[0] for item in anchored_images), default=0)) + 2
    for index, (row, view, config, items) in enumerate(group_rows, start=1):
        block_end = group_rows[index][0] if index < len(group_rows) else final_end
        camera, lens, lights = _config_parts(config)
        station = FarStation(
            station=f"工位{index}",
            view=view,
            config=config,
            items=items,
            camera=camera,
            lens=lens,
            lights=lights,
        )
        for image_row, image_column, image in sorted(anchored_images, key=lambda item: (item[0], item[1])):
            if not (row <= image_row < block_end and image_column >= image_floor):
                continue
            data = image._data()
            if not data:
                continue
            station.images.append(
                FarImage(
                    data=data,
                    extension=_image_extension(image, data),
                    caption=_image_caption(worksheet, image_row, image_column, block_end),
                )
            )
        if not station.images:
            raise OpticalFarError(f"{station.station}（{station.view}）没有读取到检测图")
        stations.append(station)
    return stations


def parse_optical_far(path: str | Path) -> OpticalFarData:
    """Parse the uniform visual FAR layout used by DCE optical workbooks."""
    source = Path(path).expanduser().resolve()
    if not source.is_file() or source.suffix.lower() != ".xlsx":
        raise OpticalFarError(f"光学 FAR 必须是存在的 .xlsx 文件：{source}")
    workbook = load_workbook(source, data_only=False, read_only=False)
    try:
        worksheet = workbook.worksheets[0]
        if "FAR" not in _cell_text(worksheet, 1, 1).upper():
            raise OpticalFarError("第一张工作表 A1 不是 FAR 分析表")
        requirements = _parse_requirements(worksheet)
        stations = _parse_stations(worksheet)
        return OpticalFarData(
            source_path=source,
            project_code=_value_after_label(worksheet, {"RFQ编号", "项目编号"}),
            production_rate=_value_after_label(worksheet, {"检测速度"}),
            special_notes=_value_after_label(worksheet, {"特别说明"}),
            requirements=requirements,
            stations=stations,
        )
    finally:
        workbook.close()


def _equipment_entry(value: str) -> tuple[str, int] | None:
    text = str(value or "").strip()
    if not text:
        return None
    count_match = _COUNT_PATTERN.search(text)
    count = int(count_match.group(1)) if count_match else 1
    name = _COUNT_PATTERN.sub("", text).strip(" ；;+/\\")
    if not name:
        return None
    return name, count


def _equipment_summary(values: list[str], *, numeric_descending: bool = False) -> str:
    totals: OrderedDict[str, int] = OrderedDict()
    for value in values:
        entry = _equipment_entry(value)
        if entry is None:
            continue
        name, count = entry
        totals[name] = totals.get(name, 0) + count
    names = list(totals)
    if numeric_descending:
        names.sort(
            key=lambda name: float(re.search(r"\d+(?:\.\d+)?", name).group())
            if re.search(r"\d+(?:\.\d+)?", name)
            else 0.0,
            reverse=True,
        )
    return "；".join(f"{name}×{totals[name]}" for name in names)


def far_equipment_summaries(data: OpticalFarData) -> tuple[str, str, str]:
    cameras = [station.camera for station in data.stations if station.camera]
    lenses = [station.lens for station in data.stations if station.lens]
    lights = [light for station in data.stations for light in station.lights]
    return (
        _equipment_summary(cameras, numeric_descending=True),
        _equipment_summary(lenses),
        _equipment_summary(lights),
    )


def _module(project: PptProject, key: str) -> ProjectModule:
    module = next(
        (item for item in project.modules if item.template_module_key == key),
        None,
    )
    if module is None:
        raise OpticalFarError(f"当前项目缺少模块：{key}")
    return module


def _template_key_for_source(module: ProjectModule, source_slide: int) -> str:
    page_template = next(
        (item for item in module.page_templates if item.source_slide == source_slide),
        None,
    )
    if page_template is None:
        raise OpticalFarError(
            f"模块“{module.name}”未登记 PPT 第 {source_slide} 页为页面模板"
        )
    return page_template.key


def _source_slide_for_project_slide(module: ProjectModule, slide: ProjectSlide) -> int:
    page_template = next(
        (item for item in module.page_templates if item.key == slide.page_template_key),
        None,
    )
    return int(page_template.source_slide) if page_template else 0


def _ensure_module_slide(module: ProjectModule, source_slide: int) -> ProjectSlide:
    existing = next(
        (
            slide
            for slide in module.slides
            if _source_slide_for_project_slide(module, slide) == source_slide
        ),
        None,
    )
    if existing is not None:
        return existing
    slide = ProjectSlide(page_template_key=_template_key_for_source(module, source_slide))
    module.slides = [slide]
    return slide


def _template_table(
    project: PptProject,
    manifest: TemplateManifest,
    slot_key: str,
) -> list[list[str]]:
    slot = next((item for item in manifest.slots if item["key"] == slot_key), None)
    if slot is None or slot["kind"] != "table":
        raise OpticalFarError(f"模板配置缺少表格 Slot：{slot_key}")
    presentation = Presentation(str(Path(project.template_path).expanduser().resolve()))
    slide = presentation.slides[int(slot["slide"]) - 1]
    shape = next((item for item in slide.shapes if item.shape_id == slot["shape_id"]), None)
    if shape is None or not getattr(shape, "has_table", False):
        raise OpticalFarError(f"模板 Slot {slot_key} 未绑定有效表格")
    return [[cell.text for cell in row.cells] for row in shape.table.rows]


def _requirement_status(mark: str) -> tuple[str, str]:
    normalized = _normalize_mark(mark)
    if "⌂" in normalized:
        return "是", "（AI）"
    if normalized == "?":
        return "待评估", ""
    if normalized.upper() == "X":
        return "否", ""
    return "是", ""


def _requirement_views(data: OpticalFarData, target: str) -> str:
    views = [
        station.view
        for station in data.stations
        if target and any(target in item or item in target for item in station.items)
    ]
    return "/".join(dict.fromkeys(view for view in views if view))


def _requirement_table(
    base_table: list[list[str]],
    data: OpticalFarData,
) -> list[list[str]]:
    if len(base_table) != 9 or any(len(row) != 5 for row in base_table):
        raise OpticalFarError("产品检测项模板表必须为 9×5")
    if len(data.requirements) > 8:
        raise OpticalFarError(
            f"FAR 含 {len(data.requirements)} 条检测项，超过模板8条容量；请先增加检测项页"
        )
    table = [list(base_table[0])]
    for index, requirement in enumerate(data.requirements, start=1):
        status, suffix = _requirement_status(requirement.mark)
        target = (requirement.target or requirement.category or "未命名检测项") + suffix
        table.append(
            [
                str(index),
                target,
                status,
                requirement.standard,
                _requirement_views(data, requirement.target),
            ]
        )
    table.extend([["", "", "", "", ""] for _ in range(9 - len(table))])
    return table


def _equipment_table(
    base_table: list[list[str]],
    data: OpticalFarData,
) -> tuple[list[list[str]], str, str, str]:
    if len(base_table) != 7 or any(len(row) != 4 for row in base_table):
        raise OpticalFarError("设备参数模板表必须为 7×4")
    camera_summary, lens_summary, light_summary = far_equipment_summaries(data)
    table = [list(row) for row in base_table]
    for row in table:
        if row[0].strip() == "生产节拍" and data.production_rate:
            row[1] = data.production_rate
        if row[2].strip() == "工业光源" and light_summary:
            row[3] = light_summary
        elif row[2].strip() == "工业相机" and camera_summary:
            row[3] = camera_summary
        elif row[2].strip() == "工业镜头" and lens_summary:
            row[3] = lens_summary
    return table, camera_summary, lens_summary, light_summary


def _persist_images(data: OpticalFarData, asset_root: Path) -> tuple[Path, list[list[Path]]]:
    digest = hashlib.sha256(data.source_path.read_bytes()).hexdigest()
    directory = Path(asset_root).expanduser().resolve() / digest
    directory.mkdir(parents=True, exist_ok=True)
    paths: list[list[Path]] = []
    for station_index, station in enumerate(data.stations, start=1):
        station_paths: list[Path] = []
        for image_index, image in enumerate(station.images, start=1):
            if image.extension == "bin":
                raise OpticalFarError(
                    f"{station.station}第 {image_index} 张图的格式无法识别"
                )
            path = directory / (
                f"station_{station_index:02d}_image_{image_index:03d}.{image.extension}"
            )
            if not path.is_file() or path.read_bytes() != image.data:
                path.write_bytes(image.data)
            station_paths.append(path)
        paths.append(station_paths)
    return directory, paths


def apply_optical_far(
    project: PptProject,
    manifest: TemplateManifest,
    data: OpticalFarData,
    asset_root: str | Path,
) -> OpticalFarApplyResult:
    """Replace FAR-driven modules in ``project`` after complete validation."""
    required_slots = {
        "far_result_title",
        "far_result_camera",
        "far_result_lens",
        "far_result_item",
        "far_result_view",
        "far_result_image",
        "far_result_note",
        "far_result_caption",
        "equipment_parameters",
        "inspection_items",
    }
    known_slots = {slot["key"] for slot in manifest.slots}
    missing_slots = sorted(required_slots - known_slots)
    if missing_slots:
        raise OpticalFarError("模板配置缺少 Slot：" + "、".join(missing_slots))
    if not data.stations or not data.image_count:
        raise OpticalFarError("FAR 没有可用检测图")

    effect_module = _module(project, "inspection_result")
    items_module = _module(project, "inspection_items")
    parameters_module = _module(project, "equipment_parameters")
    effect_template_key = _template_key_for_source(effect_module, 11)
    items_table = _requirement_table(
        _template_table(project, manifest, "inspection_items"),
        data,
    )
    equipment_table, camera_summary, lens_summary, light_summary = _equipment_table(
        _template_table(project, manifest, "equipment_parameters"),
        data,
    )
    asset_directory, image_paths = _persist_images(
        data,
        Path(asset_root),
    )

    effect_slides: list[ProjectSlide] = []
    selected_asset_paths: list[Path] = []
    page_index = 0
    for station_index, station in enumerate(data.stations):
        station_label = f"{station.station}：{station.view}"
        matched_items: list[tuple[str, int]] = []
        unmatched_items: list[str] = []
        for item in _unique_station_items(station):
            image_index = _matching_station_image_index(station, item)
            if image_index is None:
                unmatched_items.append(item)
            else:
                matched_items.append((item, image_index))

        page_specs: list[tuple[list[str], int, bool]] = []
        if unmatched_items:
            ok_image_index = _ok_station_image_index(station)
            if ok_image_index is None:
                raise OpticalFarError(
                    f"{station.station}（{station.view}）的检测项“"
                    + "、".join(unmatched_items)
                    + "”没有对应缺陷图或OK图"
                )
            page_specs.append((unmatched_items, ok_image_index, True))
        page_specs.extend(([item], image_index, False) for item, image_index in matched_items)

        for page_items, image_index, uses_ok_fallback in page_specs:
            page_index += 1
            selected_path = image_paths[station_index][image_index]
            selected_asset_paths.append(selected_path)
            full_item_text = "、".join(page_items)
            page_subject = "OK样件" if uses_ok_fallback else page_items[0]
            note_parts = [
                f"光学来源 {data.project_code or data.source_path.stem} FAR",
                station_label,
                "检测范围：" + full_item_text,
            ]
            if station.lights:
                note_parts.append("光源：" + "、".join(station.lights))
            if uses_ok_fallback:
                note_parts.append("本工位没有其它对应缺陷样图，本页展示1张OK样件图")
            effect_slides.append(
                ProjectSlide(
                    page_template_key=effect_template_key,
                    title=f"{station_label} · {page_subject}",
                    subtitle=full_item_text,
                    overrides={
                        "far_result_title": f"检测效果-{station_label}",
                        "far_result_camera": f"相机：{station.camera or '未填写'}",
                        "far_result_lens": f"镜头：{station.lens or '未填写'}",
                        "far_result_item": (
                            _ok_effect_item_text(page_items)
                            if uses_ok_fallback
                            else _effect_item_text(page_items)
                        ),
                        "far_result_view": station_label,
                        "far_result_image": str(selected_path),
                        "far_result_note": "；".join(note_parts),
                        "far_result_caption": (
                            f"图 {page_index} {station_label} OK样件图"
                            if uses_ok_fallback
                            else f"图 {page_index} {page_items[0]}检测图"
                        ),
                    },
                )
            )

    items_slide = _ensure_module_slide(items_module, 16)
    items_slide.overrides.update(
        {
            "inspection_items_title": (
                f"产品检测项说明（{data.project_code}）"
                if data.project_code
                else "产品检测项说明"
            ),
            "inspection_items": items_table,
            "inspection_accuracy_note": (
                f"注：检测标准来源 {data.source_path.name} 右上角检测项汇总。"
                + data.special_notes
            )[:260],
            "inspection_items_caption": "表 4.1 检测项说明",
        }
    )
    items_slide.title = "产品检测项说明"

    parameters_slide = _ensure_module_slide(parameters_module, 10)
    parameters_slide.overrides.update(
        {
            "equipment_parameters": equipment_table,
            "parameter_note": f"光学参数来源：{data.source_path.name}"[:120],
        }
    )

    effect_module.slides = effect_slides
    effect_module.enabled = True
    effect_module.default_add_template = effect_template_key
    items_module.enabled = True
    parameters_module.enabled = True

    project.excel_path = str(data.source_path)
    summary = (
        f"光学 FAR：{data.source_path.name}\n"
        f"项目编号：{data.project_code or '未读取'}\n"
        f"检测项：{len(data.requirements)}\n"
        f"工位：{len(data.stations)}\n"
        f"检测图：{data.image_count}\n"
        f"检测效果页：{len(effect_slides)}\n"
        f"生产节拍：{data.production_rate or '未读取'}"
    )
    source_record = next(
        (record for record in project.sources if Path(record.path) == data.source_path),
        None,
    )
    if source_record is None:
        project.sources.append(
            SourceRecord(path=str(data.source_path), kind="光学 FAR", content=summary)
        )
    else:
        source_record.kind = "光学 FAR"
        source_record.content = summary

    asset_root_path = Path(asset_root).expanduser().resolve()

    def is_generated_far_asset(asset: AssetRecord) -> bool:
        if asset.category != "检测效果图":
            return False
        try:
            Path(asset.path).expanduser().resolve().relative_to(asset_root_path)
            return True
        except ValueError:
            return False

    project.assets = [
        asset for asset in project.assets if not is_generated_far_asset(asset)
    ]
    existing_assets = {Path(asset.path).expanduser().resolve() for asset in project.assets}
    for path in dict.fromkeys(selected_asset_paths):
        resolved = path.resolve()
        if resolved in existing_assets:
            continue
        project.assets.append(AssetRecord(path=str(path), category="检测效果图"))
        existing_assets.add(resolved)

    return OpticalFarApplyResult(
        effect_pages=len(effect_slides),
        requirements=len(data.requirements),
        stations=len(data.stations),
        asset_directory=asset_directory,
        camera_summary=camera_summary,
        lens_summary=lens_summary,
        light_summary=light_summary,
    )
