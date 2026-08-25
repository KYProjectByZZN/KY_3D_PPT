"""Configuration-driven renderer for editable PowerPoint templates."""

from __future__ import annotations

import hashlib
import json
import posixpath
import tempfile
import zipfile
from collections.abc import Iterable
from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from xml.etree import ElementTree

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from .io_utils import commit_staged_output, staged_output_path
from .module_service import (
    ensure_project_modules,
    rebuild_structure_context,
    slot_specs_for_source_slide,
)
from .navigation_style import apply_navigation_style
from .project import PptProject


SUPPORTED_SLOT_KINDS = {"text", "table", "image"}


class TemplateRenderError(ValueError):
    """Raised when a template, manifest, or render value is invalid."""


@dataclass(frozen=True)
class TemplateManifest:
    """Validated template configuration used by the renderer."""

    template_filename: str
    template_sha256: str
    slide_count: int
    modules: tuple[dict[str, Any], ...]
    slots: tuple[dict[str, Any], ...]


def sha256_file(path: str | Path) -> str:
    """Return the uppercase SHA-256 digest for a file."""
    digest = hashlib.sha256()
    with Path(path).open("rb") as stream:
        for block in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest().upper()


def _read_json_object(path: Path, label: str) -> dict[str, Any]:
    try:
        content = json.loads(path.read_text(encoding="utf-8"))
    except FileNotFoundError:
        raise FileNotFoundError(f"{label}不存在：{path}") from None
    except json.JSONDecodeError as exc:
        raise TemplateRenderError(f"{label}不是有效 JSON：{path}（{exc}）") from exc
    if not isinstance(content, dict):
        raise TemplateRenderError(f"{label}顶层必须是 JSON 对象：{path}")
    return content


def load_manifest(path: str | Path) -> TemplateManifest:
    """Load and validate a version-1 template manifest."""
    manifest_path = Path(path).expanduser().resolve()
    raw = _read_json_object(manifest_path, "模板配置")

    if raw.get("schema_version") != 1:
        raise TemplateRenderError("模板配置 schema_version 必须为 1")

    template = raw.get("template")
    if not isinstance(template, dict):
        raise TemplateRenderError("模板配置缺少 template 对象")
    filename = template.get("filename")
    expected_hash = template.get("sha256")
    slide_count = template.get("slide_count")
    if not isinstance(filename, str) or not filename.strip():
        raise TemplateRenderError("template.filename 不能为空")
    if not isinstance(expected_hash, str) or len(expected_hash) != 64:
        raise TemplateRenderError("template.sha256 必须是 64 位 SHA-256")
    if not isinstance(slide_count, int) or slide_count < 1:
        raise TemplateRenderError("template.slide_count 必须是正整数")

    modules = raw.get("modules", [])
    if not isinstance(modules, list):
        raise TemplateRenderError("modules 必须是数组")
    assigned_slides: set[int] = set()
    module_keys: set[str] = set()
    for module in modules:
        if not isinstance(module, dict):
            raise TemplateRenderError("每个 module 必须是对象")
        key = module.get("key")
        slides = module.get("slides")
        if not isinstance(key, str) or not key.strip() or key in module_keys:
            raise TemplateRenderError(f"module.key 为空或重复：{key!r}")
        if not isinstance(slides, list) or not slides:
            raise TemplateRenderError(f"模块 {key} 的 slides 必须是非空数组")
        for slide_number in slides:
            if not isinstance(slide_number, int) or not 1 <= slide_number <= slide_count:
                raise TemplateRenderError(f"模块 {key} 含无效页码：{slide_number!r}")
            if slide_number in assigned_slides:
                raise TemplateRenderError(f"第 {slide_number} 页被重复分配到多个模块")
            assigned_slides.add(slide_number)
        module_keys.add(key)

    slots = raw.get("slots", [])
    if not isinstance(slots, list):
        raise TemplateRenderError("slots 必须是数组")
    slot_keys: set[str] = set()
    normalized_slots: list[dict[str, Any]] = []
    for slot in slots:
        if not isinstance(slot, dict):
            raise TemplateRenderError("每个 slot 必须是对象")
        key = slot.get("key")
        slide_number = slot.get("slide")
        shape_id = slot.get("shape_id")
        kind = slot.get("kind")
        if not isinstance(key, str) or not key.strip() or key in slot_keys:
            raise TemplateRenderError(f"slot.key 为空或重复：{key!r}")
        if not isinstance(slide_number, int) or not 1 <= slide_number <= slide_count:
            raise TemplateRenderError(f"Slot {key} 含无效页码：{slide_number!r}")
        if not isinstance(shape_id, int) or shape_id < 1:
            raise TemplateRenderError(f"Slot {key} 含无效 shape_id：{shape_id!r}")
        if kind not in SUPPORTED_SLOT_KINDS:
            raise TemplateRenderError(f"Slot {key} 的 kind 不受支持：{kind!r}")
        max_chars = slot.get("max_chars")
        if max_chars is not None and (not isinstance(max_chars, int) or max_chars < 1):
            raise TemplateRenderError(f"Slot {key} 的 max_chars 必须是正整数")
        slot_keys.add(key)
        normalized_slots.append(dict(slot))

    return TemplateManifest(
        template_filename=filename,
        template_sha256=expected_hash.upper(),
        slide_count=slide_count,
        modules=tuple(dict(module) for module in modules),
        slots=tuple(normalized_slots),
    )


def _find_shape(slide: Any, shape_id: int, *, slide_number: int, slot_key: str) -> Any:
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise TemplateRenderError(
        f"第 {slide_number} 页找不到 Shape ID {shape_id}（Slot: {slot_key}）"
    )


def _replace_text_frame(text_frame: Any, value: str) -> None:
    """Replace text while retaining the first existing run's formatting."""
    paragraphs = text_frame.paragraphs
    first_paragraph = paragraphs[0]
    if first_paragraph.runs:
        first_paragraph.runs[0].text = value
        for run in first_paragraph.runs[1:]:
            run.text = ""
    else:
        first_paragraph.add_run().text = value
    for paragraph in paragraphs[1:]:
        for run in paragraph.runs:
            run.text = ""


def _replace_text(shape: Any, value: Any, slot: dict[str, Any]) -> None:
    if not isinstance(value, str):
        raise TemplateRenderError(f"Slot {slot['key']} 的文字值必须是字符串")
    max_chars = slot.get("max_chars")
    if max_chars is not None and len(value) > max_chars:
        raise TemplateRenderError(
            f"Slot {slot['key']} 超过 {max_chars} 字限制，当前为 {len(value)} 字"
        )
    if not getattr(shape, "has_text_frame", False):
        raise TemplateRenderError(f"Slot {slot['key']} 绑定的对象不是文本对象")
    _replace_text_frame(shape.text_frame, value)


def _replace_table(shape: Any, value: Any, slot: dict[str, Any]) -> None:
    if not getattr(shape, "has_table", False):
        raise TemplateRenderError(f"Slot {slot['key']} 绑定的对象不是表格")
    if not isinstance(value, list) or not value or not all(isinstance(row, list) for row in value):
        raise TemplateRenderError(f"Slot {slot['key']} 的表格值必须是非空二维数组")

    table = shape.table
    expected_rows = len(table.rows)
    expected_columns = len(table.columns)
    if len(value) != expected_rows or any(len(row) != expected_columns for row in value):
        raise TemplateRenderError(
            f"Slot {slot['key']} 表格尺寸必须为 {expected_rows}×{expected_columns}"
        )
    for row_index, row in enumerate(value):
        for column_index, cell_value in enumerate(row):
            if not isinstance(cell_value, (str, int, float)) or isinstance(cell_value, bool):
                raise TemplateRenderError(
                    f"Slot {slot['key']} 的第 {row_index + 1} 行第 {column_index + 1} 列"
                    "必须是文字或数字"
                )
            _replace_text_frame(table.cell(row_index, column_index).text_frame, str(cell_value))


def _resolve_image_path(value: Any, *, data_directory: Path, slot_key: str) -> Path:
    if not isinstance(value, str) or not value.strip():
        raise TemplateRenderError(f"Slot {slot_key} 的图片值必须是文件路径")
    candidate = Path(value).expanduser()
    if not candidate.is_absolute():
        candidate = data_directory / candidate
    candidate = candidate.resolve()
    if not candidate.is_file():
        raise FileNotFoundError(f"Slot {slot_key} 的图片不存在：{candidate}")
    return candidate


def _resolve_slide_plan(
    manifest: TemplateManifest,
    enabled_modules: Iterable[str] | None,
    module_order: Iterable[str] | None,
) -> list[int]:
    """Return original one-based slide numbers in the requested output order."""
    default_slide_plan = list(range(1, manifest.slide_count + 1))
    if enabled_modules is None and module_order is None:
        return default_slide_plan
    if not manifest.modules:
        raise TemplateRenderError("当前模板没有 modules 配置，不能筛选或排序模块")

    module_keys = [module["key"] for module in manifest.modules]
    known_keys = set(module_keys)
    enabled_list = module_keys if enabled_modules is None else list(enabled_modules)
    if len(enabled_list) != len(set(enabled_list)):
        raise TemplateRenderError("启用模块列表包含重复项")
    unknown_enabled = sorted(set(enabled_list) - known_keys)
    if unknown_enabled:
        raise TemplateRenderError(f"启用模块不存在：{', '.join(unknown_enabled)}")
    enabled_set = set(enabled_list)

    requested_order = module_keys if module_order is None else list(module_order)
    if len(requested_order) != len(set(requested_order)):
        raise TemplateRenderError("模块顺序包含重复项")
    unknown_order = sorted(set(requested_order) - known_keys)
    if unknown_order:
        raise TemplateRenderError(f"模块顺序包含未知模块：{', '.join(unknown_order)}")
    ordered_enabled = [key for key in requested_order if key in enabled_set]
    ordered_enabled.extend(key for key in module_keys if key in enabled_set and key not in ordered_enabled)

    modules_by_key = {module["key"]: module for module in manifest.modules}
    slide_plan = [
        slide_number
        for module_key in ordered_enabled
        for slide_number in modules_by_key[module_key]["slides"]
    ]
    assigned_slides = {
        slide_number for module in manifest.modules for slide_number in module["slides"]
    }
    slide_plan.extend(
        slide_number for slide_number in default_slide_plan if slide_number not in assigned_slides
    )
    if not slide_plan:
        raise TemplateRenderError("至少需要启用一个包含页面的模块")
    return slide_plan


def _apply_slide_plan(presentation: Any, slide_plan: list[int]) -> None:
    """Filter and reorder slides without copying their relationship graphs."""
    slide_id_list = presentation.slides._sldIdLst
    original_slide_ids = list(slide_id_list)
    selected = set(slide_plan)
    for slide_number, slide_id in enumerate(original_slide_ids, start=1):
        if slide_number not in selected:
            presentation.part.drop_rel(slide_id.rId)
    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)
    for slide_number in slide_plan:
        slide_id_list.append(original_slide_ids[slide_number - 1])


def _replace_image(
    slide: Any,
    shape: Any,
    value: Any,
    slot: dict[str, Any],
    *,
    data_directory: Path,
) -> None:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        raise TemplateRenderError(f"Slot {slot['key']} 绑定的对象不是图片")
    image_path = _resolve_image_path(value, data_directory=data_directory, slot_key=slot["key"])

    parent = shape._element.getparent()
    original_index = parent.index(shape._element)
    geometry = (shape.left, shape.top, shape.width, shape.height)
    crop = (shape.crop_left, shape.crop_top, shape.crop_right, shape.crop_bottom)
    original_name = shape.name
    parent.remove(shape._element)

    replacement = slide.shapes.add_picture(str(image_path), *geometry)
    replacement.crop_left, replacement.crop_top, replacement.crop_right, replacement.crop_bottom = crop
    replacement.name = original_name
    parent.remove(replacement._element)
    parent.insert(original_index, replacement._element)


def _rewrite_relationship_ids(element: Any, relationship_ids: dict[str, str]) -> None:
    """Rewrite relationship references in a copied Open XML element tree."""
    if not relationship_ids:
        return
    for child in element.iter():
        for attribute, value in list(child.attrib.items()):
            replacement = relationship_ids.get(value)
            if replacement:
                child.set(attribute, replacement)


def _remove_slide_tag_metadata(element: Any) -> None:
    """Drop non-visual per-slide tags from a cloned slide.

    PowerPoint treats a ``ppt/tags/tag*.xml`` part as owned by exactly one
    slide.  Reusing the same tag part across multiple clones produces a PPTX
    that python-pptx can reopen but desktop PowerPoint refuses to open.  The
    tags in this template are WPS editing metadata and have no visual effect,
    so cloned pages intentionally omit them.
    """
    for child in list(element.iter()):
        if child.tag.endswith("}custDataLst"):
            parent = child.getparent()
            if parent is not None:
                parent.remove(child)


def _relate_like(destination_part: Any, relationship: Any) -> str:
    if relationship.is_external:
        return destination_part.relate_to(
            relationship.target_ref,
            relationship.reltype,
            is_external=True,
        )
    return destination_part.relate_to(relationship.target_part, relationship.reltype)


def _relationship_id(destination_part: Any, relationship_type: str) -> str:
    matches = [
        relationship.rId
        for relationship in destination_part.rels.values()
        if relationship.reltype == relationship_type
    ]
    if len(matches) != 1:
        raise TemplateRenderError(
            f"复制幻灯片时关系 {relationship_type.rsplit('/', 1)[-1]} 数量异常"
        )
    return matches[0]


def _clone_notes_slide(source_slide: Any, destination_slide: Any) -> None:
    if not source_slide.has_notes_slide:
        return
    source_notes = source_slide.notes_slide
    destination_notes = destination_slide.notes_slide
    relationship_ids: dict[str, str] = {}
    for relationship in source_notes.part.rels.values():
        if relationship.reltype in {RT.NOTES_MASTER, RT.SLIDE}:
            relationship_ids[relationship.rId] = _relationship_id(
                destination_notes.part, relationship.reltype
            )
        else:
            relationship_ids[relationship.rId] = _relate_like(
                destination_notes.part, relationship
            )
    copied_element = deepcopy(source_notes._element)
    _rewrite_relationship_ids(copied_element, relationship_ids)
    destination_notes._element = copied_element
    destination_notes.part._element = copied_element
    destination_notes.__dict__.pop("shapes", None)
    destination_notes.__dict__.pop("placeholders", None)


def _clone_slide(presentation: Any, source_slide: Any) -> Any:
    """Clone a slide and its relationship graph inside one presentation package."""
    destination_slide = presentation.slides.add_slide(source_slide.slide_layout)
    relationship_ids: dict[str, str] = {}
    source_notes_relationship_id = ""
    for relationship in source_slide.part.rels.values():
        if relationship.reltype == RT.NOTES_SLIDE:
            source_notes_relationship_id = relationship.rId
            continue
        if relationship.reltype == RT.TAGS:
            continue
        if relationship.reltype == RT.SLIDE_LAYOUT:
            relationship_ids[relationship.rId] = _relationship_id(
                destination_slide.part, RT.SLIDE_LAYOUT
            )
        else:
            relationship_ids[relationship.rId] = _relate_like(
                destination_slide.part, relationship
            )

    copied_element = deepcopy(source_slide._element)
    _remove_slide_tag_metadata(copied_element)
    _rewrite_relationship_ids(copied_element, relationship_ids)
    destination_slide._element = copied_element
    destination_slide.part._element = copied_element
    # ``add_slide()`` populates and caches layout placeholders before we replace
    # the XML tree. Drop those cached collections so callers see the copied tree.
    destination_slide.__dict__.pop("shapes", None)
    destination_slide.__dict__.pop("placeholders", None)

    if source_slide.has_notes_slide:
        _clone_notes_slide(source_slide, destination_slide)
        if source_notes_relationship_id:
            notes_id = _relationship_id(destination_slide.part, RT.NOTES_SLIDE)
            _rewrite_relationship_ids(
                destination_slide._element,
                {source_notes_relationship_id: notes_id},
            )
    return destination_slide


def _remove_original_slides(
    presentation: Any,
    original_count: int,
    replacements: dict[Any, Any],
) -> None:
    # Some real-world templates contain non-standard reverse relationships from
    # layouts or tag parts to a slide. Redirect those to the first cloned
    # instance so an otherwise-unreferenced original slide is not serialized as
    # an orphan package part.
    original_parts = {
        slide.part for slide in list(presentation.slides)[:original_count]
    }
    fallback_part = next(iter(replacements.values()), None)
    for part in list(presentation.part.package.iter_parts()):
        for relationship in part.rels.values():
            if relationship.is_external:
                continue
            replacement = replacements.get(relationship.target_part)
            if (
                replacement is None
                and fallback_part is not None
                and relationship.target_part in original_parts
            ):
                replacement = fallback_part
            if replacement is not None:
                relationship._target = replacement
    slide_id_list = presentation.slides._sldIdLst
    original_slide_ids = list(slide_id_list)[:original_count]
    for slide_id in original_slide_ids:
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _remove_unreferenced_slide_parts(path: Path) -> None:
    """Remove slide XML serialized only through a template's reverse relation.

    A few WPS/PowerPoint-authored templates contain layout-to-slide reverse
    relationships. ``python-pptx`` can serialize the old slide part even after
    it is removed from ``p:sldIdLst``. The presentation never references that
    page, so remove only those unreferenced slide members and their content-type
    declarations after saving.
    """
    relationship_namespace = "http://schemas.openxmlformats.org/package/2006/relationships"
    content_type_namespace = "http://schemas.openxmlformats.org/package/2006/content-types"
    with zipfile.ZipFile(path) as package:
        relationship_root = ElementTree.fromstring(
            package.read("ppt/_rels/presentation.xml.rels")
        )
        referenced = {
            posixpath.normpath(posixpath.join("ppt", relationship.attrib["Target"]))
            for relationship in relationship_root.findall(
                f"{{{relationship_namespace}}}Relationship"
            )
            if relationship.attrib.get("Type") == RT.SLIDE
        }
        slide_entries = {
            name
            for name in package.namelist()
            if name.startswith("ppt/slides/slide")
            and name.endswith(".xml")
            and "/_rels/" not in name
        }
        unreferenced = slide_entries - referenced
        if not unreferenced:
            return
        content_types = ElementTree.fromstring(package.read("[Content_Types].xml"))
        for override in list(content_types):
            part_name = override.attrib.get("PartName", "").lstrip("/")
            if part_name in unreferenced:
                content_types.remove(override)
        ElementTree.register_namespace("", content_type_namespace)
        content_types_data = ElementTree.tostring(
            content_types,
            encoding="utf-8",
            xml_declaration=True,
        )
        related_members = {
            f"ppt/slides/_rels/{Path(name).name}.rels"
            for name in unreferenced
        }
        members_to_write = [
            (
                member,
                content_types_data
                if member.filename == "[Content_Types].xml"
                else package.read(member.filename),
            )
            for member in package.infolist()
            if member.filename not in unreferenced
            and member.filename not in related_members
        ]

    with tempfile.NamedTemporaryFile(
        prefix=f"{path.stem}_clean_",
        suffix=".pptx",
        dir=path.parent,
        delete=False,
    ) as temp_stream:
        temporary_path = Path(temp_stream.name)
    try:
        with zipfile.ZipFile(temporary_path, "w") as destination:
            for member, data in members_to_write:
                destination.writestr(member, data)
        temporary_path.replace(path)
    finally:
        if temporary_path.exists():
            temporary_path.unlink()


def validate_pptx_package(path: str | Path, *, expected_slides: int | None = None) -> None:
    """Validate the basic ZIP/Open XML structure of a generated PPTX file."""
    package_path = Path(path)
    if package_path.read_bytes()[:4] != b"PK\x03\x04":
        raise TemplateRenderError(f"输出文件没有有效的 ZIP 签名：{package_path}")
    try:
        with zipfile.ZipFile(package_path) as package:
            bad_member = package.testzip()
            if bad_member:
                raise TemplateRenderError(f"PPTX ZIP 成员损坏：{bad_member}")
            names = set(package.namelist())
    except zipfile.BadZipFile as exc:
        raise TemplateRenderError(f"输出文件不是有效 PPTX ZIP：{package_path}") from exc

    required = {"[Content_Types].xml", "ppt/presentation.xml"}
    missing = sorted(required - names)
    if missing:
        raise TemplateRenderError(f"PPTX 缺少必要文件：{', '.join(missing)}")
    if expected_slides is not None:
        slide_entries = {
            name
            for name in names
            if name.startswith("ppt/slides/slide")
            and name.endswith(".xml")
            and "/_rels/" not in name
        }
        if len(slide_entries) != expected_slides:
            raise TemplateRenderError(
                f"PPTX 页面 XML 数量应为 {expected_slides}，实际为 {len(slide_entries)}"
            )


def render_template(
    template_path: str | Path,
    manifest_path: str | Path,
    data_path: str | Path,
    output_path: str | Path,
    *,
    overwrite: bool = False,
    enabled_modules: Iterable[str] | None = None,
    module_order: Iterable[str] | None = None,
) -> Path:
    """Render configured values into a copied PPTX template."""
    template_file = Path(template_path).expanduser().resolve()
    manifest_file = Path(manifest_path).expanduser().resolve()
    data_file = Path(data_path).expanduser().resolve()
    destination = Path(output_path).expanduser().resolve()

    if template_file.suffix.lower() != ".pptx" or not template_file.is_file():
        raise FileNotFoundError(f"PPTX 模板不存在：{template_file}")
    if destination.suffix.lower() != ".pptx":
        raise TemplateRenderError("输出文件必须使用 .pptx 扩展名")
    if destination == template_file:
        raise TemplateRenderError("输出路径不能与原模板相同")
    if destination.exists() and not overwrite:
        raise FileExistsError(f"目标文件已经存在：{destination}")

    manifest = load_manifest(manifest_file)
    if template_file.name != manifest.template_filename:
        raise TemplateRenderError(
            f"模板文件名与配置不符：应为 {manifest.template_filename}，实际为 {template_file.name}"
        )
    source_hash = sha256_file(template_file)
    if source_hash != manifest.template_sha256:
        raise TemplateRenderError(
            f"模板 SHA-256 与配置不符：应为 {manifest.template_sha256}，实际为 {source_hash}"
        )

    slide_plan = _resolve_slide_plan(manifest, enabled_modules, module_order)
    selected_slides = set(slide_plan)

    values = _read_json_object(data_file, "渲染数据")
    known_keys = {slot["key"] for slot in manifest.slots}
    unknown_keys = sorted(set(values) - known_keys)
    if unknown_keys:
        raise TemplateRenderError(f"渲染数据包含未配置 Slot：{', '.join(unknown_keys)}")
    missing_required = [
        slot["key"]
        for slot in manifest.slots
        if slot.get("required", False)
        and slot["slide"] in selected_slides
        and (
            slot["key"] not in values
            or values[slot["key"]] is None
            or (slot["kind"] == "text" and not str(values[slot["key"]]).strip())
        )
    ]
    if missing_required:
        raise TemplateRenderError(f"渲染数据缺少必填 Slot：{', '.join(missing_required)}")

    presentation = Presentation(template_file)
    if len(presentation.slides) != manifest.slide_count:
        raise TemplateRenderError(
            f"模板页数与配置不符：应为 {manifest.slide_count}，实际为 {len(presentation.slides)}"
        )

    resolved_slots: list[tuple[dict[str, Any], Any, Any]] = []
    for slot in manifest.slots:
        if slot["slide"] not in selected_slides:
            continue
        slide = presentation.slides[slot["slide"] - 1]
        shape = _find_shape(
            slide,
            slot["shape_id"],
            slide_number=slot["slide"],
            slot_key=slot["key"],
        )
        resolved_slots.append((slot, slide, shape))

    for slot, slide, shape in resolved_slots:
        if slot["key"] not in values:
            continue
        value = values[slot["key"]]
        if slot["kind"] == "text":
            _replace_text(shape, value, slot)
        elif slot["kind"] == "table":
            _replace_table(shape, value, slot)
        else:
            _replace_image(slide, shape, value, slot, data_directory=data_file.parent)

    _apply_slide_plan(presentation, slide_plan)
    with staged_output_path(destination) as staged:
        presentation.save(staged)
        validate_pptx_package(staged, expected_slides=len(slide_plan))
        if sha256_file(template_file) != source_hash:
            raise TemplateRenderError("原模板在生成过程中发生了变化")
        commit_staged_output(staged, destination)
    return destination


def render_project(
    project: PptProject,
    output_path: str | Path | None = None,
    *,
    overwrite: bool = False,
) -> Path:
    """Render schema-v2 module and slide instances into an editable PPTX."""
    template_file = Path(project.template_path).expanduser().resolve()
    manifest_file = Path(project.manifest_path).expanduser().resolve()
    destination_value = output_path or project.output_path
    if not destination_value:
        raise TemplateRenderError("项目没有设置输出路径")
    destination = Path(destination_value).expanduser().resolve()

    if template_file.suffix.lower() != ".pptx" or not template_file.is_file():
        raise FileNotFoundError(f"PPTX 模板不存在：{template_file}")
    if destination.suffix.lower() != ".pptx":
        raise TemplateRenderError("输出文件必须使用 .pptx 扩展名")
    if destination == template_file:
        raise TemplateRenderError("输出路径不能与原模板相同")
    if destination.exists() and not overwrite:
        raise FileExistsError(f"目标文件已经存在：{destination}")

    manifest = load_manifest(manifest_file)
    if template_file.name != manifest.template_filename:
        raise TemplateRenderError(
            f"模板文件名与配置不符：应为 {manifest.template_filename}，实际为 {template_file.name}"
        )
    source_hash = sha256_file(template_file)
    if source_hash != manifest.template_sha256:
        raise TemplateRenderError(
            f"模板 SHA-256 与配置不符：应为 {manifest.template_sha256}，实际为 {source_hash}"
        )

    ensure_project_modules(project, manifest)
    contexts = rebuild_structure_context(project, manifest)
    if not contexts:
        raise TemplateRenderError("至少需要启用一个包含页面的模块")

    presentation = Presentation(template_file)
    if len(presentation.slides) != manifest.slide_count:
        raise TemplateRenderError(
            f"模板页数与配置不符：应为 {manifest.slide_count}，实际为 {len(presentation.slides)}"
        )
    source_slides = list(presentation.slides)
    data_directory = manifest_file.parent
    source_replacements: dict[Any, Any] = {}

    for context in contexts:
        _append_project_context(
            presentation,
            source_slides,
            context,
            manifest,
            data_directory,
            source_replacements,
        )

    _remove_original_slides(
        presentation,
        manifest.slide_count,
        source_replacements,
    )
    apply_navigation_style(
        presentation,
        project.presentation_style,
        [
            project.presentation_style.navigation_index_for(
                context.template_module_key
            )
            for context in contexts
        ],
    )
    with staged_output_path(destination) as staged:
        presentation.save(staged)
        _remove_unreferenced_slide_parts(staged)
        validate_pptx_package(staged, expected_slides=len(contexts))
        reopened = Presentation(staged)
        if len(reopened.slides) != len(contexts):
            raise TemplateRenderError("生成后的 PPTX 重开页数不一致")
        if sha256_file(template_file) != source_hash:
            raise TemplateRenderError("原模板在生成过程中发生了变化")
        commit_staged_output(staged, destination)
    return destination


def _append_project_context(
    presentation: Any,
    source_slides: list[Any],
    context: Any,
    manifest: TemplateManifest,
    data_directory: Path,
    source_replacements: dict[Any, Any],
) -> Any:
    """Clone and fill one project context inside an open presentation."""
    if not 1 <= context.source_slide <= len(source_slides):
        raise TemplateRenderError(
            f"页面 {context.slide_title} 引用了无效模板页：{context.source_slide}"
        )
    source_slide = source_slides[context.source_slide - 1]
    slide = _clone_slide(presentation, source_slide)
    source_replacements.setdefault(source_slide.part, slide.part)
    if context.remove_shape_ids:
        for shape in list(slide.shapes):
            if shape.shape_id in set(context.remove_shape_ids):
                shape._element.getparent().remove(shape._element)
    slot_specs = slot_specs_for_source_slide(manifest, context.source_slide)
    missing_required = [
        slot["key"]
        for slot in slot_specs
        if slot.get("required", False)
        and (
            slot["key"] not in context.values
            or context.values[slot["key"]] is None
            or (
                slot["kind"] == "text"
                and not str(context.values[slot["key"]]).strip()
            )
        )
    ]
    if missing_required:
        raise TemplateRenderError(
            f"模块“{context.module_title}”页面“{context.slide_title}”缺少必填 Slot："
            + ", ".join(missing_required)
        )

    for slot in slot_specs:
        key = slot["key"]
        if key not in context.values:
            continue
        shape = _find_shape(
            slide,
            slot["shape_id"],
            slide_number=context.source_slide,
            slot_key=key,
        )
        value = context.values[key]
        if slot["kind"] == "text":
            _replace_text(shape, value, slot)
        elif slot["kind"] == "table":
            _replace_table(shape, value, slot)
        else:
            _replace_image(
                slide,
                shape,
                value,
                slot,
                data_directory=data_directory,
            )

    for system_key, shape_id in context.system_slots.items():
        value_key = "page_number_text" if system_key == "page_number" else system_key
        value = context.values.get(value_key, "")
        shape = _find_shape(
            slide,
            shape_id,
            slide_number=context.source_slide,
            slot_key=system_key,
        )
        _replace_text(
            shape,
            str(value),
            {"key": system_key, "kind": "text"},
        )
    return slide


def render_project_page(
    project: PptProject,
    module_id: str,
    slide_id: str,
    output_path: str | Path,
    *,
    overwrite: bool = False,
) -> Path:
    """Render one selected page while retaining its full-project context."""
    template_file = Path(project.template_path).expanduser().resolve()
    manifest_file = Path(project.manifest_path).expanduser().resolve()
    destination = Path(output_path).expanduser().resolve()
    if template_file.suffix.lower() != ".pptx" or not template_file.is_file():
        raise FileNotFoundError(f"PPTX 模板不存在：{template_file}")
    if destination.suffix.lower() != ".pptx":
        raise TemplateRenderError("输出文件必须使用 .pptx 扩展名")
    if destination == template_file:
        raise TemplateRenderError("输出路径不能与原模板相同")
    if destination.exists() and not overwrite:
        raise FileExistsError(f"目标文件已经存在：{destination}")

    manifest = load_manifest(manifest_file)
    if template_file.name != manifest.template_filename:
        raise TemplateRenderError(
            f"模板文件名与配置不符：应为 {manifest.template_filename}，实际为 {template_file.name}"
        )
    source_hash = sha256_file(template_file)
    if source_hash != manifest.template_sha256:
        raise TemplateRenderError(
            f"模板 SHA-256 与配置不符：应为 {manifest.template_sha256}，实际为 {source_hash}"
        )

    ensure_project_modules(project, manifest)
    context = next(
        (
            item
            for item in rebuild_structure_context(project, manifest)
            if item.module_id == module_id and item.slide_id == slide_id
        ),
        None,
    )
    if context is None:
        raise TemplateRenderError("无法在最终 PPT 结构中定位当前页面")

    presentation = Presentation(template_file)
    if len(presentation.slides) != manifest.slide_count:
        raise TemplateRenderError(
            f"模板页数与配置不符：应为 {manifest.slide_count}，实际为 {len(presentation.slides)}"
        )
    source_slides = list(presentation.slides)
    source_replacements: dict[Any, Any] = {}
    _append_project_context(
        presentation,
        source_slides,
        context,
        manifest,
        manifest_file.parent,
        source_replacements,
    )
    _remove_original_slides(
        presentation,
        manifest.slide_count,
        source_replacements,
    )
    apply_navigation_style(
        presentation,
        project.presentation_style,
        [
            project.presentation_style.navigation_index_for(
                context.template_module_key
            )
        ],
    )
    with staged_output_path(destination) as staged:
        presentation.save(staged)
        _remove_unreferenced_slide_parts(staged)
        validate_pptx_package(staged, expected_slides=1)
        if len(Presentation(staged).slides) != 1:
            raise TemplateRenderError("单页预览 PPTX 重开页数不一致")
        if sha256_file(template_file) != source_hash:
            raise TemplateRenderError("原模板在生成过程中发生了变化")
        commit_staged_output(staged, destination)
    return destination
