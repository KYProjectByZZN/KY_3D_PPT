from __future__ import annotations

import tempfile
import unittest
import zipfile
from pathlib import Path

from pptx import Presentation

from ppt_generator import build_presentation


class BuilderTests(unittest.TestCase):
    def test_generates_valid_two_slide_presentation(self) -> None:
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            output = Path(temp_dir) / "sample.pptx"
            result = build_presentation("项目汇报", "KY Project", ["背景", "方案"], output)

            self.assertEqual(result, output.resolve())
            self.assertEqual(output.read_bytes()[:4], b"PK\x03\x04")

            with zipfile.ZipFile(output) as package:
                names = set(package.namelist())
            self.assertIn("[Content_Types].xml", names)
            self.assertIn("ppt/presentation.xml", names)
            self.assertIn("ppt/slides/slide1.xml", names)
            self.assertIn("ppt/slides/slide2.xml", names)

            presentation = Presentation(output)
            self.assertEqual(len(presentation.slides), 2)
            self.assertEqual(presentation.slides[0].shapes.title.text, "项目汇报")

    def test_existing_output_requires_explicit_overwrite(self) -> None:
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            output = Path(temp_dir) / "sample.pptx"
            build_presentation("第一次", "", [], output)

            with self.assertRaises(FileExistsError):
                build_presentation("第二次", "", [], output)

            build_presentation("第二次", "", [], output, overwrite=True)
            presentation = Presentation(output)
            self.assertEqual(presentation.slides[0].shapes.title.text, "第二次")

    def test_rejects_invalid_output_extension(self) -> None:
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            with self.assertRaises(ValueError):
                build_presentation("项目", "", [], Path(temp_dir) / "sample.pdf")


if __name__ == "__main__":
    unittest.main()
