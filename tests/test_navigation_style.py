from __future__ import annotations

import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches

from ppt_generator.module_service import ensure_project_modules
from ppt_generator.navigation_style import (
    NAVIGATION_ITEMS,
)
from ppt_generator.project import NavigationItem, PptProject, PresentationStyle
from ppt_generator.template_renderer import (
    load_manifest,
    render_project,
    sha256_file,
)


PROJECT_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = PROJECT_ROOT / "templates" / "冲压筒形壳体检测方案NAT6704_v2.pptx"
MANIFEST = PROJECT_ROOT / "templates" / "NAT6704_v2.template.json"
DATA = PROJECT_ROOT / "examples" / "NAT6704_v2_test_data.json"


class NavigationStyleTests(unittest.TestCase):
    def test_project_renderer_applies_navigation_style(self) -> None:
        source_hash = sha256_file(TEMPLATE)
        manifest = load_manifest(MANIFEST)
        project = PptProject(
            template_path=str(TEMPLATE),
            manifest_path=str(MANIFEST),
            values=json.loads(DATA.read_text(encoding="utf-8")),
            presentation_style=PresentationStyle(
                navigation_height=0.60,
                navigation_background="#DDEEFF",
            ),
        )
        ensure_project_modules(project, manifest)
        target = next(
            module
            for module in project.modules
            if any(page.source_slide == 5 for page in module.page_templates)
        )
        for module in project.modules:
            module.enabled = module.id == target.id

        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            output = Path(temp_dir) / "navigation_style.pptx"
            render_project(project, output)
            reopened = Presentation(output)

        self.assertEqual(sha256_file(TEMPLATE), source_hash)
        self.assertEqual(len(reopened.slides), len(target.slides))
        layout = reopened.slides[0].slide_layout
        background = next(
            shape
            for shape in layout.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and shape.width > Inches(12.0)
            and shape.top < Inches(0.95)
        )
        self.assertAlmostEqual(background.height / Inches(1), 0.60, places=2)
        self.assertEqual(str(background.fill.fore_color.rgb), "FFFFFF")
        background_effect = background.element.spPr.find(qn("a:effectLst"))
        self.assertTrue(
            background_effect is None
            or background_effect.find(qn("a:outerShdw")) is None
        )

        navigation_text = [
            shape
            for shape in reopened.slides[0].shapes
            if shape.name.startswith("KY_NAV_TEXT_")
        ]
        self.assertEqual(len(navigation_text), 5)
        self.assertEqual(
            tuple(shape.text.strip() for shape in navigation_text),
            NAVIGATION_ITEMS,
        )
        for shape in navigation_text:
            runs = [
                run
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
            ]
            self.assertTrue(runs)
            self.assertTrue(all(run.font.bold for run in runs))
            self.assertTrue(
                all(
                    run.font.size is not None
                    and abs(
                        run.font.size.pt
                        - project.presentation_style.resolved_navigation_font_size()
                    )
                    < 0.01
                    for run in runs
                )
            )
        active_text = next(
            shape for shape in navigation_text if shape.text.strip() == "设备设计"
        )
        active_runs = [
            run
            for paragraph in active_text.text_frame.paragraphs
            for run in paragraph.runs
        ]
        self.assertTrue(
            all(str(run.font.color.rgb) == "C90000" for run in active_runs)
        )
        active_background = next(
            shape
            for shape in reopened.slides[0].shapes
            if shape.name == "KY_NAV_ACTIVE_BACKGROUND"
        )
        active_bar = next(
            shape
            for shape in reopened.slides[0].shapes
            if shape.name == "KY_NAV_ACTIVE"
        )
        baseline = next(
            shape
            for shape in reopened.slides[0].shapes
            if shape.name == "KY_NAV_BASELINE"
        )
        self.assertEqual(str(active_background.fill.fore_color.rgb), "DDEEFF")
        self.assertEqual(str(baseline.fill.fore_color.rgb), "D3D9DE")
        self.assertEqual(baseline.left, 0)
        self.assertLessEqual(
            abs(baseline.left + baseline.width - reopened.slide_width),
            1,
        )
        self.assertEqual(baseline.top, active_bar.top)
        self.assertEqual(baseline.height, active_bar.height)
        baseline_effect = baseline.element.spPr.find(qn("a:effectLst"))
        self.assertTrue(
            baseline_effect is None
            or baseline_effect.find(qn("a:outerShdw")) is None
        )
        self.assertEqual(active_background.left, active_text.left)
        self.assertEqual(active_background.width, active_text.width)
        self.assertEqual(active_bar.left, active_text.left)
        self.assertEqual(active_bar.width, active_text.width)
        self.assertFalse(
            any(
                shape.name.startswith("KY_NAV_INACTIVE_")
                for shape in reopened.slides[0].shapes
            )
        )
        self.assertEqual(navigation_text[0].left, 0)
        for previous, current in zip(navigation_text, navigation_text[1:]):
            self.assertLessEqual(abs(previous.left + previous.width - current.left), 1)
        self.assertLessEqual(
            abs(navigation_text[-1].left + navigation_text[-1].width - Inches(11.55)),
            1,
        )
        self.assertFalse(
            any(
                getattr(shape, "has_text_frame", False)
                and shape.text.strip() in {"企业简介", "设备布局", "模块介绍", "供货范围"}
                for shape in layout.shapes
            )
        )

    def test_renderer_supports_added_sixth_navigation_item(self) -> None:
        manifest = load_manifest(MANIFEST)
        project = PptProject(
            template_path=str(TEMPLATE),
            manifest_path=str(MANIFEST),
            values=json.loads(DATA.read_text(encoding="utf-8")),
        )
        project.presentation_style.navigation_items.append(
            NavigationItem("交付范围")
        )
        project.presentation_style.navigation_font_size = 13.0
        ensure_project_modules(project, manifest)
        target = next(
            module
            for module in project.modules
            if module.template_module_key == "vision_system"
        )
        for module in project.modules:
            module.enabled = module.id == target.id

        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            output = Path(temp_dir) / "six_navigation_items.pptx"
            render_project(project, output)
            reopened = Presentation(output)

        navigation_text = [
            shape
            for shape in reopened.slides[0].shapes
            if shape.name.startswith("KY_NAV_TEXT_")
        ]
        self.assertEqual(
            [shape.text.strip() for shape in navigation_text],
            ["公司简介", "工艺分析", "设备设计", "检测效果", "系统介绍", "交付范围"],
        )
        self.assertEqual(
            len(
                [
                    shape
                    for shape in reopened.slides[0].shapes
                    if shape.name.startswith("KY_NAV_SEPARATOR_")
                ]
            ),
            5,
        )
        active_text = next(
            shape for shape in navigation_text if shape.text.strip() == "系统介绍"
        )
        self.assertTrue(
            all(
                run.font.size is not None and abs(run.font.size.pt - 13.0) < 0.01
                for shape in navigation_text
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
            )
        )
        self.assertTrue(
            all(
                str(run.font.color.rgb) == "C90000"
                for paragraph in active_text.text_frame.paragraphs
                for run in paragraph.runs
            )
        )


if __name__ == "__main__":
    unittest.main()
