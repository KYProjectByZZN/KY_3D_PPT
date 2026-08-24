from __future__ import annotations

import tempfile
import unittest
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from ppt_generator import PptProject, ensure_project_modules, load_manifest
from ppt_generator.optical_far import (
    apply_optical_far,
    far_equipment_summaries,
    parse_optical_far,
)
from ppt_generator.template_renderer import _clone_slide


PROJECT_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = PROJECT_ROOT / "templates" / "冲压筒形壳体检测方案NAT6704_v2.pptx"
MANIFEST = PROJECT_ROOT / "templates" / "NAT6704_v2.template.json"
FAR = PROJECT_ROOT / "templates" / "光学资料" / "NAT6801FAR(8.5).xlsx"


class OpticalFarTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.data = parse_optical_far(FAR)

    def test_parses_real_far_requirements_stations_and_images(self) -> None:
        data = self.data
        self.assertEqual(data.project_code, "NAT6801")
        self.assertEqual(data.production_rate, "40-50pcs/min")
        self.assertEqual(len(data.requirements), 5)
        self.assertEqual(len(data.stations), 3)
        self.assertEqual(data.image_count, 21)
        self.assertEqual(data.requirements[0].target, "翘曲")
        self.assertEqual(
            data.requirements[0].standard,
            "长边单边局部区域高度超差0.2mm",
        )
        self.assertEqual(data.stations[1].images[0].caption, "指纹")
        self.assertEqual(data.stations[2].images[-1].caption, "翘曲")
        self.assertTrue(data.stations[0].images[0].data.startswith(b"\x89PNG"))

    def test_summarizes_camera_lens_and_light(self) -> None:
        self.assertEqual(
            far_equipment_summaries(self.data),
            (
                "1200万×1；500万×2",
                "FA镜头×2；远心镜头×1",
                "同轴光×1；背光×3",
            ),
        )

    def test_duplicate_effect_pages_drop_single_owner_tag_metadata(self) -> None:
        presentation = Presentation(TEMPLATE)
        source_slide = presentation.slides[10]
        clones = [
            _clone_slide(presentation, source_slide),
            _clone_slide(presentation, source_slide),
        ]
        for clone in clones:
            self.assertFalse(
                any(rel.reltype == RT.TAGS for rel in clone.part.rels.values())
            )
            self.assertFalse(
                any(
                    element.tag.endswith("}custDataLst")
                    for element in clone._element.iter()
                )
            )

    def test_applies_far_to_project_without_duplicate_assets(self) -> None:
        manifest = load_manifest(MANIFEST)
        project = PptProject(
            template_path=str(TEMPLATE),
            manifest_path=str(MANIFEST),
        )
        ensure_project_modules(project, manifest)
        data = deepcopy(self.data)
        data.stations[0].items.append("指纹")
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            result = apply_optical_far(
                project,
                manifest,
                data,
                Path(temp_dir) / "far_assets",
            )
            self.assertEqual(result.effect_pages, 6)
            effect_module = next(
                item
                for item in project.modules
                if item.template_module_key == "inspection_result"
            )
            self.assertEqual(len(effect_module.slides), 6)
            first = effect_module.slides[0].overrides
            second = effect_module.slides[1].overrides
            last = effect_module.slides[-1].overrides
            self.assertEqual(first["far_result_camera"], "相机：500万*1")
            self.assertEqual(
                first["far_result_item"],
                "OK样件图（\n指纹、镀层不良、生锈、磨伤）",
            )
            self.assertEqual(first["far_result_view"], "工位1：俯视")
            self.assertIn("OK样件", first["far_result_caption"])
            self.assertIn("镀层不良（发黑、发黄）", first["far_result_note"])
            self.assertIn("没有其它对应缺陷样图", first["far_result_note"])
            self.assertEqual(second["far_result_item"], "检测项：指纹")
            self.assertEqual(second["far_result_view"], "工位2：仰视")
            self.assertNotIn("OK样件", second["far_result_caption"])
            self.assertEqual(last["far_result_item"], "检测项：翘曲")
            self.assertEqual(last["far_result_view"], "工位3：侧视")
            self.assertTrue(Path(first["far_result_image"]).is_file())
            self.assertEqual(
                len(
                    {
                        slide.overrides["far_result_image"]
                        for slide in effect_module.slides
                    }
                ),
                6,
            )
            self.assertEqual(len(project.assets), 6)
            self.assertEqual(
                [
                    Path(slide.overrides["far_result_image"]).name
                    for slide in effect_module.slides
                ],
                [
                    "station_01_image_001.png",
                    "station_02_image_001.png",
                    "station_02_image_002.png",
                    "station_02_image_004.png",
                    "station_02_image_008.png",
                    "station_03_image_003.png",
                ],
            )
            page_keys = [
                (
                    slide.overrides["far_result_view"],
                    slide.overrides["far_result_item"],
                )
                for slide in effect_module.slides
            ]
            self.assertEqual(len(page_keys), len(set(page_keys)))
            self.assertEqual(
                sum(len(module.slides) for module in project.modules if module.enabled),
                24,
            )

            items_module = next(
                item
                for item in project.modules
                if item.template_module_key == "inspection_items"
            )
            items_table = items_module.slides[0].overrides["inspection_items"]
            self.assertEqual((len(items_table), len(items_table[0])), (9, 5))
            self.assertEqual(items_table[1][1:4], ["翘曲", "是", "长边单边局部区域高度超差0.2mm"])
            self.assertEqual(items_table[2][1], "指纹（AI）")
            self.assertEqual(items_table[2][4], "俯视/仰视")

            parameters_module = next(
                item
                for item in project.modules
                if item.template_module_key == "equipment_parameters"
            )
            parameter_table = parameters_module.slides[0].overrides[
                "equipment_parameters"
            ]
            self.assertEqual(parameter_table[0], ["生产节拍", "40-50pcs/min", "工业光源", "同轴光×1；背光×3"])
            self.assertEqual(parameter_table[1][3], "1200万×1；500万×2")
            self.assertEqual(parameter_table[2][3], "FA镜头×2；远心镜头×1")
            self.assertEqual(parameter_table[3][1], "Win10-64bit")

            apply_optical_far(
                project,
                manifest,
                data,
                Path(temp_dir) / "far_assets",
            )
            self.assertEqual(len(effect_module.slides), 6)
            self.assertEqual(len(project.assets), 6)


if __name__ == "__main__":
    unittest.main()
