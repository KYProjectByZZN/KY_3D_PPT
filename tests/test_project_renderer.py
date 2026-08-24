from __future__ import annotations

import json
import tempfile
import unittest
import zipfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches

from ppt_generator.module_service import duplicate_module, ensure_project_modules
from ppt_generator.project import PptProject
from ppt_generator.template_renderer import (
    load_manifest,
    render_project,
    render_project_page,
    sha256_file,
)


class ProjectRendererTests(unittest.TestCase):
    def _fixture(self, directory: Path) -> tuple[Path, Path]:
        image_path = directory / "source.png"
        Image.new("RGB", (160, 100), "#2f6fad").save(image_path)
        template = directory / "module_template.pptx"
        presentation = Presentation()
        shape_ids: list[int] = []
        for index in range(2):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            title = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.4), Inches(5), Inches(0.8)
            )
            title.text = f"模板标题{index + 1}"
            title.text_frame.paragraphs[0].runs[0].font.bold = True
            shape_ids.append(title.shape_id)
            if index == 0:
                slide.shapes.add_picture(
                    str(image_path), Inches(0.5), Inches(1.5), Inches(3), Inches(2)
                )
                run = title.text_frame.paragraphs[0].runs[0]
                run.hyperlink.address = "https://example.com/module"
                slide.notes_slide.notes_text_frame.text = "模块页面备注"
        # Reproduce the reverse layout-to-slide relationship found in the real
        # NAT6704 file. Without cleanup it serializes an unreferenced old slide.
        presentation.slide_layouts[6].part.relate_to(
            presentation.slides[0].part,
            RT.SLIDE,
        )
        presentation.save(template)

        manifest = directory / "module_template.json"
        manifest.write_text(
            json.dumps(
                {
                    "schema_version": 1,
                    "template": {
                        "filename": template.name,
                        "sha256": sha256_file(template),
                        "slide_count": 2,
                    },
                    "modules": [
                        {"key": "optical", "name": "光学方案", "slides": [1, 2]}
                    ],
                    "slots": [
                        {
                            "key": "page_one_title",
                            "slide": 1,
                            "shape_id": shape_ids[0],
                            "kind": "text",
                            "required": True,
                        },
                        {
                            "key": "page_two_title",
                            "slide": 2,
                            "shape_id": shape_ids[1],
                            "kind": "text",
                            "required": True,
                        },
                    ],
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )
        return template, manifest

    def test_renders_independent_module_copies_and_relationships(self) -> None:
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            directory = Path(temp_dir)
            template, manifest_path = self._fixture(directory)
            source_hash = sha256_file(template)
            manifest = load_manifest(manifest_path)
            project = PptProject(
                template_path=str(template),
                manifest_path=str(manifest_path),
            )
            ensure_project_modules(project, manifest)
            source = project.modules[0]
            source.slides[0].overrides["page_one_title"] = "工位A方案"
            source.slides[1].overrides["page_two_title"] = "工位A参数"
            copied = duplicate_module(project, source.id)
            copied.slides[0].overrides["page_one_title"] = "工位B方案"
            copied.slides[1].overrides["page_two_title"] = "工位B参数"
            output = directory / "result.pptx"

            render_project(project, output)

            self.assertEqual(sha256_file(template), source_hash)
            self.assertEqual(output.read_bytes()[:4], b"PK\x03\x04")
            with zipfile.ZipFile(output) as package:
                self.assertIsNone(package.testzip())
                self.assertEqual(
                    len(
                        [
                            name
                            for name in package.namelist()
                            if name.startswith("ppt/slides/slide")
                            and name.endswith(".xml")
                            and "/_rels/" not in name
                        ]
                    ),
                    4,
                )

            result = Presentation(output)
            self.assertEqual(len(result.slides), 4)
            titles = [
                next(
                    shape.text
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False) and shape.text
                )
                for slide in result.slides
            ]
            self.assertEqual(
                titles,
                ["工位A方案", "工位A参数", "工位B方案", "工位B参数"],
            )
            self.assertTrue(
                any(
                    shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    for shape in result.slides[0].shapes
                )
            )
            self.assertIn("模块页面备注", result.slides[0].notes_slide.notes_text_frame.text)
            self.assertEqual(
                result.slides[0]
                .shapes[0]
                .text_frame.paragraphs[0]
                .runs[0]
                .hyperlink.address,
                "https://example.com/module",
            )

            preview_output = directory / "single_page_preview.pptx"
            render_project_page(
                project,
                copied.id,
                copied.slides[1].id,
                preview_output,
            )
            self.assertEqual(sha256_file(template), source_hash)
            with zipfile.ZipFile(preview_output) as package:
                self.assertIsNone(package.testzip())
                self.assertEqual(
                    len(
                        [
                            name
                            for name in package.namelist()
                            if name.startswith("ppt/slides/slide")
                            and name.endswith(".xml")
                            and "/_rels/" not in name
                        ]
                    ),
                    1,
                )
            preview = Presentation(preview_output)
            self.assertEqual(len(preview.slides), 1)
            self.assertIn("工位B参数", preview.slides[0].shapes[0].text)


if __name__ == "__main__":
    unittest.main()
