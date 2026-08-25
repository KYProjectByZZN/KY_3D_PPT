from __future__ import annotations

import base64
from copy import deepcopy
import json
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from ppt_generator.no_cad_scheme import NoCadSchemeService
from ppt_generator.module_service import ensure_project_modules
from ppt_generator.project import DeviceModule, PptProject, load_project, save_project
from ppt_generator.scheme_application import (
    import_no_cad_scene,
    sync_no_cad_scene_to_ppt,
)
from ppt_generator.scheme_service import SchemeError, materialize_equipment_scheme
from ppt_generator.template_renderer import load_manifest, render_project


PROJECT_ROOT = Path(__file__).resolve().parents[1]
PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
)


class SchemeApplicationTests(unittest.TestCase):
    def test_project_identity_mismatch_rejects_sync_without_mutation(self) -> None:
        service = NoCadSchemeService()
        scene = service.create_minimum_scene()
        scene.project_name = "其它项目"
        project = PptProject(project_name="NAT6704")
        manifest = load_manifest(
            PROJECT_ROOT / "templates" / "NAT6704_v2.template.json"
        )
        ensure_project_modules(project, manifest)
        before_scheme = deepcopy(project.equipment_scheme.to_dict())
        before_modules = deepcopy([value.to_dict() for value in project.modules])

        with self.assertRaisesRegex(SchemeError, "项目名称"):
            sync_no_cad_scene_to_ppt(project, scene, manifest)

        self.assertEqual(project.equipment_scheme.to_dict(), before_scheme)
        self.assertEqual([value.to_dict() for value in project.modules], before_modules)

    def test_missing_images_imports_formal_scheme_without_touching_ppt_modules(self) -> None:
        service = NoCadSchemeService()
        scene = service.create_minimum_scene()
        project = PptProject(project_name=scene.project_name)
        manifest = load_manifest(
            PROJECT_ROOT / "templates" / "NAT6704_v2.template.json"
        )
        ensure_project_modules(project, manifest)
        before_modules = deepcopy([value.to_dict() for value in project.modules])

        synced = sync_no_cad_scene_to_ppt(project, scene, manifest)

        self.assertIsNone(synced.materialization)
        self.assertFalse(synced.ppt_updated)
        self.assertIn("整机方案", synced.pending_image_names)
        self.assertEqual(
            len(project.equipment_scheme.equipment_modules),
            len(scene.nodes),
        )
        self.assertEqual([value.to_dict() for value in project.modules], before_modules)

    def test_all_confirmed_images_sync_formal_scheme_and_ppt_in_one_action(self) -> None:
        service = NoCadSchemeService()
        scene = service.create_minimum_scene()
        project = PptProject(
            project_name=scene.project_name,
            template_path=str(
                PROJECT_ROOT / "templates" / "冲压筒形壳体检测方案NAT6704_v2.pptx"
            ),
            manifest_path=str(
                PROJECT_ROOT / "templates" / "NAT6704_v2.template.json"
            ),
            values=json.loads(
                (
                    PROJECT_ROOT / "examples" / "NAT6704_v2_test_data.json"
                ).read_text(encoding="utf-8")
            ),
        )
        manifest = load_manifest(
            PROJECT_ROOT / "templates" / "NAT6704_v2.template.json"
        )
        ensure_project_modules(project, manifest)

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            result = service.evaluate(scene)
            overview = root / "overview.png"
            overview.write_bytes(PNG_1X1)
            service.bind_accepted_image(
                scene,
                "overview",
                str(overview),
                {"targetHash": result.visual_target("overview").target_hash},
            )
            expected: list[str] = []
            for index, node in enumerate(scene.nodes, start=1):
                image = root / f"module-{index}.png"
                image.write_bytes(PNG_1X1)
                expected.append(str(image.resolve()))
                service.bind_accepted_image(
                    scene,
                    node.node_id,
                    str(image),
                    {"targetHash": result.visual_target(node.node_id).target_hash},
                )

            synced = sync_no_cad_scene_to_ppt(project, scene, manifest)
            output = root / "project-visual-sync.pptx"
            project.output_path = str(output)
            render_project(project)
            self.assertEqual(output.read_bytes()[:4], b"PK\x03\x04")
            self.assertGreater(len(Presentation(output).slides), len(scene.nodes))

        self.assertTrue(synced.ppt_updated)
        self.assertIsNotNone(synced.materialization)
        self.assertEqual(synced.pending_image_names, ())
        overview_module = next(
            value for value in project.modules
            if value.template_module_key == "equipment_overview"
        )
        equipment_module = next(
            value for value in project.modules
            if value.template_module_key == "equipment_module"
        )
        self.assertEqual(
            overview_module.slides[0].overrides["equipment_image"],
            str(overview.resolve()),
        )
        actual = [
            next(
                value for key, value in slide.overrides.items()
                if key.startswith("equipment_module_image_")
            )
            for slide in equipment_module.slides
        ]
        self.assertEqual(actual, expected)

    def test_imports_overview_and_each_module_with_prompt_structure_and_provenance(self) -> None:
        service = NoCadSchemeService()
        scene = service.create_minimum_scene()
        result = service.evaluate(scene)
        project = PptProject()
        manual = DeviceModule(name="人工补充模块")
        project.equipment_scheme.equipment_modules.append(manual)

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            overview = root / "overview.png"
            module = root / "module.png"
            overview.write_bytes(b"overview")
            module.write_bytes(b"module")
            service.bind_accepted_image(
                scene,
                "overview",
                str(overview),
                {"targetHash": result.visual_target("overview").target_hash},
            )
            service.bind_accepted_image(
                scene,
                scene.nodes[1].node_id,
                str(module),
                {"targetHash": result.visual_target(scene.nodes[1].node_id).target_hash},
            )

            imported = import_no_cad_scene(project, scene)
            imported_ids = {
                value.source_scene_node_id: value.id
                for value in project.equipment_scheme.equipment_modules
                if value.source_scene_node_id
            }
            second = import_no_cad_scene(project, scene)

        scheme = project.equipment_scheme
        self.assertEqual(imported.flow_nodes, 3)
        self.assertEqual(imported.equipment_modules, 3)
        self.assertEqual(imported.image_targets, 4)
        self.assertEqual(imported.pending_images, 2)
        self.assertEqual(second.pending_images, 2)
        self.assertIn(manual, scheme.equipment_modules)
        self.assertEqual(
            imported_ids,
            {
                value.source_scene_node_id: value.id
                for value in scheme.equipment_modules
                if value.source_scene_node_id
            },
        )
        vision = next(
            value
            for value in scheme.equipment_modules
            if value.source_scene_node_id == scene.nodes[1].node_id
        )
        self.assertTrue(vision.structure_definition["components"])
        self.assertIn("Authoritative module structure JSON", vision.image_prompt)
        self.assertEqual(vision.image_path, str(module))
        self.assertEqual(scheme.overview_image, str(overview))
        self.assertIn("moduleStructures", scheme.overview_structure)
        self.assertEqual(len(project.assets), 2)

        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "module-visuals.kyppt.json"
            save_project(project, path)
            restored = load_project(path)
        restored_vision = next(
            value
            for value in restored.equipment_scheme.equipment_modules
            if value.source_scene_node_id == scene.nodes[1].node_id
        )
        self.assertEqual(restored_vision.structure_definition, vision.structure_definition)
        self.assertEqual(restored_vision.image_prompt, vision.image_prompt)
        self.assertEqual(restored_vision.image_provenance, vision.image_provenance)
        self.assertEqual(restored.assets[0].metadata["source"], "no-cad-ai")

    def test_all_accepted_targets_materialize_to_overview_and_module_image_slots(self) -> None:
        service = NoCadSchemeService()
        scene = service.create_minimum_scene()
        added = service.add_module(scene, "stop_position", index=1)
        service.auto_layout(scene)
        project = PptProject()
        manifest = load_manifest(
            PROJECT_ROOT / "templates" / "NAT6704_v2.template.json"
        )
        ensure_project_modules(project, manifest)

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            overview = root / "overview.png"
            overview.write_bytes(PNG_1X1)
            result = service.evaluate(scene)
            service.bind_accepted_image(
                scene,
                "overview",
                str(overview),
                {"targetHash": result.visual_target("overview").target_hash},
            )
            expected_module_images: list[str] = []
            for index, node in enumerate(scene.nodes, start=1):
                image = root / f"module-{index}.png"
                image.write_bytes(PNG_1X1)
                expected_module_images.append(str(image.resolve()))
                service.bind_accepted_image(
                    scene,
                    node.node_id,
                    str(image),
                    {"targetHash": result.visual_target(node.node_id).target_hash},
                )

            imported = import_no_cad_scene(project, scene)
            materialized = materialize_equipment_scheme(project, manifest)

        overview_module = next(
            value
            for value in project.modules
            if value.template_module_key == "equipment_overview"
        )
        equipment_module = next(
            value
            for value in project.modules
            if value.template_module_key == "equipment_module"
        )
        self.assertEqual(imported.pending_images, 0)
        self.assertIn(
            added.node_id,
            {
                value.source_scene_node_id
                for value in project.equipment_scheme.equipment_modules
            },
        )
        self.assertEqual(materialized.equipment_pages, len(scene.nodes))
        self.assertEqual(
            overview_module.slides[0].overrides["equipment_image"],
            str(overview.resolve()),
        )
        actual_module_images = [
            next(
                value
                for key, value in slide.overrides.items()
                if key.startswith("equipment_module_image_")
            )
            for slide in equipment_module.slides
        ]
        self.assertEqual(actual_module_images, expected_module_images)


if __name__ == "__main__":
    unittest.main()
