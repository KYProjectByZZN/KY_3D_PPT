from __future__ import annotations

import json
import tempfile
import unittest
import zipfile
from pathlib import Path
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from ppt_generator import TemplateRenderError, render_template, sha256_file


class TemplateRendererTests(unittest.TestCase):
    def _make_fixture(self, directory: Path) -> tuple[Path, Path, Path, bytes]:
        original_image = directory / "original.png"
        replacement_image = directory / "replacement.png"
        Image.new("RGB", (120, 80), "#336699").save(original_image)
        Image.new("RGB", (120, 80), "#CC6633").save(replacement_image)

        source = directory / "fixture.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        text_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.6))
        text_shape.text = "原始标题"
        text_shape.text_frame.paragraphs[0].runs[0].font.bold = True
        table_shape = slide.shapes.add_table(2, 2, Inches(0.5), Inches(1.5), Inches(5), Inches(1.5))
        table_shape.table.cell(0, 0).text = "原始A"
        table_shape.table.cell(0, 1).text = "原始B"
        table_shape.table.cell(1, 0).text = "原始C"
        table_shape.table.cell(1, 1).text = "原始D"
        picture_shape = slide.shapes.add_picture(
            str(original_image), Inches(6), Inches(1.5), Inches(3), Inches(2)
        )
        presentation.save(source)
        source_bytes = source.read_bytes()

        manifest = directory / "manifest.json"
        manifest.write_text(
            json.dumps(
                {
                    "schema_version": 1,
                    "template": {
                        "filename": source.name,
                        "sha256": sha256_file(source),
                        "slide_count": 1,
                    },
                    "modules": [{"key": "test", "name": "测试", "slides": [1]}],
                    "slots": [
                        {
                            "key": "title",
                            "slide": 1,
                            "shape_id": text_shape.shape_id,
                            "kind": "text",
                            "required": True,
                            "max_chars": 30,
                        },
                        {
                            "key": "table",
                            "slide": 1,
                            "shape_id": table_shape.shape_id,
                            "kind": "table",
                            "required": True,
                        },
                        {
                            "key": "image",
                            "slide": 1,
                            "shape_id": picture_shape.shape_id,
                            "kind": "image",
                            "required": True,
                        },
                    ],
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )

        data = directory / "data.json"
        data.write_text(
            json.dumps(
                {
                    "title": "替换后的标题",
                    "table": [["新A", "新B"], ["新C", 12.5]],
                    "image": replacement_image.name,
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        return source, manifest, data, source_bytes

    def test_replaces_text_table_and_image_without_touching_source(self) -> None:
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            directory = Path(temp_dir)
            source, manifest, data, source_bytes = self._make_fixture(directory)
            output = directory / "result.pptx"

            result = render_template(source, manifest, data, output)

            self.assertEqual(result, output.resolve())
            self.assertEqual(source.read_bytes(), source_bytes)
            self.assertEqual(output.read_bytes()[:4], b"PK\x03\x04")
            with zipfile.ZipFile(output) as package:
                self.assertIsNone(package.testzip())
                self.assertIn("ppt/slides/slide1.xml", package.namelist())

            presentation = Presentation(output)
            slide = presentation.slides[0]
            text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
            self.assertEqual(text_shapes[0].text, "替换后的标题")
            self.assertTrue(text_shapes[0].text_frame.paragraphs[0].runs[0].font.bold)

            table = next(shape.table for shape in slide.shapes if getattr(shape, "has_table", False))
            self.assertEqual(table.cell(0, 0).text, "新A")
            self.assertEqual(table.cell(1, 1).text, "12.5")

            picture = next(
                shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            )
            self.assertEqual(picture.image.blob, (directory / "replacement.png").read_bytes())

    def test_rejects_existing_output_without_overwrite(self) -> None:
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            directory = Path(temp_dir)
            source, manifest, data, _ = self._make_fixture(directory)
            output = directory / "result.pptx"
            output.write_bytes(b"existing")

            with self.assertRaises(FileExistsError):
                render_template(source, manifest, data, output)

    def test_validation_failure_preserves_existing_output(self) -> None:
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            directory = Path(temp_dir)
            source, manifest, data, _ = self._make_fixture(directory)
            output = directory / "result.pptx"
            original = b"existing-valid-output-placeholder"
            output.write_bytes(original)

            with patch(
                "ppt_generator.template_renderer.validate_pptx_package",
                side_effect=TemplateRenderError("模拟校验失败"),
            ):
                with self.assertRaisesRegex(TemplateRenderError, "模拟校验失败"):
                    render_template(source, manifest, data, output, overwrite=True)

            self.assertEqual(output.read_bytes(), original)
            self.assertEqual(list(directory.glob(".result.tmp-*.pptx")), [])

    def test_rejects_manifest_hash_mismatch(self) -> None:
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            directory = Path(temp_dir)
            source, manifest, data, _ = self._make_fixture(directory)
            raw_manifest = json.loads(manifest.read_text(encoding="utf-8"))
            raw_manifest["template"]["sha256"] = "0" * 64
            manifest.write_text(json.dumps(raw_manifest), encoding="utf-8")

            with self.assertRaisesRegex(TemplateRenderError, "SHA-256"):
                render_template(source, manifest, data, directory / "result.pptx")

    def test_rejects_unknown_data_key(self) -> None:
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            directory = Path(temp_dir)
            source, manifest, data, _ = self._make_fixture(directory)
            raw_data = json.loads(data.read_text(encoding="utf-8"))
            raw_data["typo"] = "不会被静默忽略"
            data.write_text(json.dumps(raw_data), encoding="utf-8")

            with self.assertRaisesRegex(TemplateRenderError, "未配置 Slot"):
                render_template(source, manifest, data, directory / "result.pptx")

    def test_filters_and_reorders_modules(self) -> None:
        with tempfile.TemporaryDirectory(dir=Path(__file__).parent) as temp_dir:
            directory = Path(temp_dir)
            source = directory / "modules.pptx"
            presentation = Presentation()
            for title in ("第一页", "第二页", "第三页"):
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(4), Inches(0.6)
                ).text = title
            presentation.save(source)

            manifest = directory / "manifest.json"
            manifest.write_text(
                json.dumps(
                    {
                        "schema_version": 1,
                        "template": {
                            "filename": source.name,
                            "sha256": sha256_file(source),
                            "slide_count": 3,
                        },
                        "modules": [
                            {"key": "a", "name": "A", "slides": [1]},
                            {"key": "b", "name": "B", "slides": [2]},
                            {"key": "c", "name": "C", "slides": [3]},
                        ],
                        "slots": [],
                    }
                ),
                encoding="utf-8",
            )
            data = directory / "data.json"
            data.write_text("{}", encoding="utf-8")

            output = directory / "result.pptx"
            render_template(
                source,
                manifest,
                data,
                output,
                enabled_modules=["c", "a"],
                module_order=["c", "a"],
            )

            result = Presentation(output)
            self.assertEqual(len(result.slides), 2)
            self.assertEqual(result.slides[0].shapes[0].text, "第三页")
            self.assertEqual(result.slides[1].shapes[0].text, "第一页")
            with zipfile.ZipFile(output) as package:
                slide_entries = [
                    name
                    for name in package.namelist()
                    if name.startswith("ppt/slides/slide")
                    and name.endswith(".xml")
                    and "/_rels/" not in name
                ]
            self.assertEqual(len(slide_entries), 2)


if __name__ == "__main__":
    unittest.main()
