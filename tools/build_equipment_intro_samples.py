"""Build five editable industrial equipment-introduction sample pages."""

from __future__ import annotations

import argparse
import hashlib
import re
import sys
import zipfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from ppt_generator.navigation_style import (
    NAVIGATION_ITEMS,
    apply_navigation_style,
)
from ppt_generator.optical_far import OpticalFarData, parse_optical_far
from ppt_generator.project import PresentationStyle


TEMPLATE = PROJECT_ROOT / "templates" / "冲压筒形壳体检测方案NAT6704_v2.pptx"
FAR_SOURCE = PROJECT_ROOT / "templates" / "光学资料" / "NAT6801FAR(8.5).xlsx"
DEFAULT_OUTPUT = PROJECT_ROOT / "output" / "设备介绍_工业版五页样本_v2.pptx"

FONT_NAME = "Microsoft YaHei"

COLORS = {
    "red": RGBColor(0xC9, 0x00, 0x00),
    "text": RGBColor(0x22, 0x2A, 0x33),
    "muted": RGBColor(0x66, 0x72, 0x7E),
    "blue": RGBColor(0x31, 0x5F, 0x7D),
    "light_blue": RGBColor(0xE9, 0xF0, 0xF5),
    "light": RGBColor(0xF4, 0xF6, 0xF7),
    "line": RGBColor(0xD3, 0xD9, 0xDE),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}


@dataclass(frozen=True)
class NavigationConfig:
    height: float = 0.52
    background: str = "#FFFFFF"
    minimum_height: float = 0.42
    maximum_height: float = 0.72

    def validate(self) -> None:
        PresentationStyle(
            navigation_height=self.height,
            navigation_background=self.background,
        ).validate()

    @property
    def title_top(self) -> float:
        return self.height + 0.18

    @property
    def title_rule_top(self) -> float:
        return self.height + 0.70

    @property
    def body_top(self) -> float:
        return self.height + 0.91


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest().upper()


def _set_typeface(run) -> None:
    run.font.name = FONT_NAME
    properties = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = properties.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            properties.append(node)
        node.set("typeface", FONT_NAME)


def _set_run(
    run,
    *,
    size: float,
    color: RGBColor,
    bold: bool = False,
) -> None:
    _set_typeface(run)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: float,
    color: RGBColor = COLORS["text"],
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.02,
):
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = vertical
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = 0
    paragraph.space_after = 0
    paragraph.line_spacing = 1.05
    run = paragraph.add_run()
    run.text = text
    _set_run(run, size=size, color=color, bold=bold)
    return shape


def _add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor = COLORS["white"],
    line: RGBColor = COLORS["line"],
    line_width: float = 0.7,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def _add_line(
    slide,
    left: float,
    top: float,
    width: float,
    *,
    color: RGBColor = COLORS["line"],
    line_width: float = 0.8,
):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(left),
        Inches(top),
        Inches(left + width),
        Inches(top),
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(line_width)
    return shape


def _picture_blob(
    presentation: Presentation,
    slide_number: int,
    shape_id: int,
) -> bytes:
    shape = next(
        item
        for item in presentation.slides[slide_number - 1].shapes
        if item.shape_id == shape_id
    )
    return shape.image.blob


def _add_picture_contain(
    slide,
    blob: bytes,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    background: RGBColor = COLORS["white"],
    border: bool = True,
    padding: float = 0.06,
):
    if border:
        _add_rect(slide, left, top, width, height, fill=background)
    with Image.open(BytesIO(blob)) as image:
        image_width, image_height = image.size
    available_width = max(0.1, width - padding * 2)
    available_height = max(0.1, height - padding * 2)
    scale = min(available_width / image_width, available_height / image_height)
    display_width = image_width * scale
    display_height = image_height * scale
    return slide.shapes.add_picture(
        BytesIO(blob),
        Inches(left + (width - display_width) / 2),
        Inches(top + (height - display_height) / 2),
        Inches(display_width),
        Inches(display_height),
    )


def _add_caption(slide, text: str, left: float, top: float, width: float) -> None:
    _add_text(
        slide,
        text,
        left,
        top,
        width,
        0.25,
        size=8.8,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
    )


def _remove_other_slides(presentation: Presentation, keep_slides: list[object]) -> None:
    keep_parts = {slide.part for slide in keep_slides}
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        if presentation.part.related_part(slide_id.rId) in keep_parts:
            continue
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape.element)


def _find_layout(presentation: Presentation, name: str):
    return next(layout for layout in presentation.slide_layouts if layout.name == name)


def _configure_navigation(
    presentation: Presentation,
    config: NavigationConfig,
    active_indices: list[int | None] | None = None,
) -> None:
    config.validate()
    apply_navigation_style(
        presentation,
        PresentationStyle(
            navigation_height=config.height,
            navigation_background=config.background,
        ),
        active_indices,
    )


def _add_page_title(slide, title: str, section: str, config: NavigationConfig) -> None:
    _add_rect(
        slide,
        0.48,
        config.title_top + 0.03,
        0.055,
        0.34,
        fill=COLORS["red"],
        line=COLORS["red"],
        line_width=0,
    )
    _add_text(
        slide,
        title,
        0.61,
        config.title_top,
        8.0,
        0.42,
        size=20.0,
        bold=True,
        vertical=MSO_ANCHOR.MIDDLE,
    )
    _add_text(
        slide,
        section,
        9.68,
        config.title_top + 0.06,
        3.10,
        0.28,
        size=8.0,
        color=COLORS["muted"],
        bold=True,
        align=PP_ALIGN.RIGHT,
        vertical=MSO_ANCHOR.MIDDLE,
    )
    _add_line(slide, 0.48, config.title_rule_top, 12.35, color=COLORS["line"])


def _add_footer(slide, page_number: int, note: str = "") -> None:
    if note:
        _add_text(
            slide,
            note,
            0.55,
            7.06,
            9.80,
            0.22,
            size=7.0,
            color=COLORS["muted"],
            vertical=MSO_ANCHOR.MIDDLE,
        )
    _add_line(slide, 11.94, 7.15, 0.36, color=COLORS["red"], line_width=1.2)
    _add_text(
        slide,
        f"{page_number:02d}",
        12.34,
        7.02,
        0.47,
        0.25,
        size=8.0,
        color=COLORS["muted"],
        align=PP_ALIGN.RIGHT,
        vertical=MSO_ANCHOR.MIDDLE,
    )


def _add_metric_grid(slide, metrics: list[tuple[str, str]], left: float, top: float) -> None:
    width = 3.88
    row_height = 0.55
    column_width = width / 2
    _add_text(slide, "关键指标", left, top, width, 0.30, size=9.2, color=COLORS["blue"], bold=True)
    top += 0.37
    _add_line(slide, left, top, width, color=COLORS["blue"], line_width=1.0)
    for index, (label, value) in enumerate(metrics):
        column = index % 2
        row = index // 2
        x = left + column * column_width
        y = top + row * row_height
        if column == 1:
            divider = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x), Inches(y + 0.10), Inches(x), Inches(y + row_height - 0.10)
            )
            divider.line.color.rgb = COLORS["line"]
            divider.line.width = Pt(0.6)
        _add_text(slide, label, x + 0.10, y + 0.08, 0.72, 0.22, size=7.7, color=COLORS["muted"])
        _add_text(
            slide,
            value,
            x + 0.74,
            y + 0.07,
            column_width - 0.84,
            0.25,
            size=8.6,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
        _add_line(slide, x, y + row_height, column_width, color=COLORS["line"], line_width=0.5)


def _add_number_marker(slide, number: str, left: float, top: float) -> None:
    _add_rect(
        slide,
        left,
        top,
        0.38,
        0.30,
        fill=COLORS["blue"],
        line=COLORS["blue"],
        line_width=0,
    )
    _add_text(
        slide,
        number,
        left,
        top,
        0.38,
        0.30,
        size=7.5,
        color=COLORS["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def _add_module_row(
    slide,
    number: str,
    name: str,
    description: str,
    left: float,
    top: float,
    width: float,
) -> None:
    _add_number_marker(slide, number, left, top + 0.08)
    _add_text(slide, name, left + 0.55, top, width - 0.55, 0.28, size=10.0, bold=True)
    _add_text(slide, description, left + 0.55, top + 0.32, width - 0.58, 0.34, size=7.8, color=COLORS["muted"])
    _add_line(slide, left, top + 0.79, width, color=COLORS["line"], line_width=0.6)


def _add_key_value_rows(
    slide,
    title: str,
    rows: list[tuple[str, str]],
    left: float,
    top: float,
    width: float,
    *,
    row_height: float = 0.42,
) -> float:
    _add_text(slide, title, left, top, width, 0.28, size=9.2, color=COLORS["blue"], bold=True)
    top += 0.34
    _add_line(slide, left, top, width, color=COLORS["blue"], line_width=1.0)
    for label, value in rows:
        _add_text(slide, label, left + 0.06, top + 0.08, 1.06, 0.22, size=7.8, color=COLORS["muted"])
        _add_text(
            slide,
            value,
            left + 1.16,
            top + 0.07,
            width - 1.22,
            0.24,
            size=8.3,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
        top += row_height
        _add_line(slide, left, top, width, color=COLORS["line"], line_width=0.5)
    return top


def _add_action_chain(
    slide,
    steps: list[str],
    top: float,
    *,
    left: float = 0.58,
    width: float = 12.06,
) -> None:
    _add_text(slide, "动作链", left, top, 0.75, 0.30, size=8.7, color=COLORS["blue"], bold=True)
    steps_left = left + 0.84
    available = width - 0.84
    cell_width = available / len(steps)
    for index, step in enumerate(steps):
        x = steps_left + index * cell_width
        _add_text(slide, f"{index + 1:02d}", x, top, 0.35, 0.28, size=7.7, color=COLORS["blue"], bold=True)
        _add_text(slide, step, x + 0.38, top, cell_width - 0.55, 0.28, size=8.0, bold=True)
        if index < len(steps) - 1:
            _add_text(
                slide,
                "→",
                x + cell_width - 0.19,
                top - 0.02,
                0.18,
                0.30,
                size=10.0,
                color=COLORS["muted"],
                align=PP_ALIGN.CENTER,
            )
    _add_line(slide, steps_left, top + 0.36, available, color=COLORS["line"], line_width=0.6)


def _build_overview(slide, images: dict[str, bytes], config: NavigationConfig) -> None:
    _add_page_title(slide, "3.1  设备总体方案", "设备布局 / EQUIPMENT LAYOUT", config)
    top = config.body_top
    _add_picture_contain(slide, images["overview"], 0.55, top, 7.72, 4.82, background=COLORS["light"])
    _add_caption(slide, "图 3.1  设备总体图", 0.55, top + 4.87, 7.72)

    _add_text(slide, "多工位视觉检测设备", 8.61, top + 0.02, 4.05, 0.45, size=17.0, bold=True)
    _add_text(
        slide,
        "面向筒形壳体多视角外观检测，集成输送定位、姿态切换、视觉成像与分类下料。",
        8.61,
        top + 0.57,
        3.92,
        0.68,
        size=9.2,
        color=COLORS["muted"],
    )
    _add_line(slide, 8.61, top + 1.39, 3.92, color=COLORS["line"])
    metrics = [
        ("检测视角", "3个"),
        ("功能模块", "4组"),
        ("生产节拍", "40–50 pcs/min"),
        ("产品流向", "左进右出"),
        ("结果输出", "OK / NG"),
        ("控制方式", "PLC联动"),
    ]
    _add_metric_grid(slide, metrics, 8.61, top + 1.68)
    _add_text(
        slide,
        "样本数据，正式项目以机械、电气和光学工程师确认值为准。",
        8.61,
        top + 4.44,
        3.92,
        0.35,
        size=7.2,
        color=COLORS["red"],
        align=PP_ALIGN.RIGHT,
    )
    _add_footer(slide, 1)


def _build_module_map(slide, images: dict[str, bytes], config: NavigationConfig) -> None:
    _add_page_title(slide, "3.2  功能模块分区", "设备布局 / EQUIPMENT LAYOUT", config)
    top = config.body_top
    _add_picture_contain(slide, images["module_map"], 0.55, top, 8.18, 4.82)
    for number, left, marker_top in (
        ("01", 1.18, top + 2.58),
        ("02", 3.72, top + 2.54),
        ("03", 5.68, top + 2.48),
        ("04", 7.66, top + 1.50),
    ):
        _add_number_marker(slide, number, left, marker_top)
    _add_caption(slide, "图 3.2  设备功能模块分区图", 0.55, top + 4.87, 8.18)

    _add_text(slide, "功能模块", 9.02, top + 0.02, 3.62, 0.30, size=9.2, color=COLORS["blue"], bold=True)
    _add_line(slide, 9.02, top + 0.39, 3.62, color=COLORS["blue"], line_width=1.0)
    modules = [
        ("01", "上料与定位", "产品导入、缓存与初始定位"),
        ("02", "搬运与翻转", "工位间移载与产品姿态切换"),
        ("03", "视觉检测", "多视角成像与结果输出"),
        ("04", "分类下料", "按判定结果分流OK/NG产品"),
    ]
    for index, row in enumerate(modules):
        _add_module_row(slide, *row, 9.02, top + 0.56 + index * 1.02, 3.62)
    _add_text(
        slide,
        "编号位置为样本；正式方案由工程师在分区图中确认。",
        9.02,
        top + 4.73,
        3.62,
        0.28,
        size=7.0,
        color=COLORS["muted"],
        align=PP_ALIGN.RIGHT,
    )
    _add_footer(slide, 2)


def _build_module_intro(slide, images: dict[str, bytes], config: NavigationConfig) -> None:
    _add_page_title(slide, "3.2.2  搬运与翻转模块", "模块介绍 / MODULE DESCRIPTION", config)
    top = config.body_top
    _add_picture_contain(slide, images["handling"], 0.55, top, 7.30, 3.90, background=COLORS["light"])
    _add_caption(slide, "图 3.3  搬运与翻转模块结构图", 0.55, top + 3.95, 7.30)

    _add_text(slide, "功能定位", 8.16, top, 4.45, 0.28, size=9.2, color=COLORS["blue"], bold=True)
    _add_text(
        slide,
        "完成产品在检测工位之间的稳定移载与姿态切换，为多视角成像提供重复定位。",
        8.16,
        top + 0.36,
        4.34,
        0.58,
        size=9.0,
    )
    _add_line(slide, 8.16, top + 1.04, 4.40, color=COLORS["line"])
    _add_text(slide, "核心功能", 8.16, top + 1.23, 4.40, 0.27, size=9.2, color=COLORS["blue"], bold=True)
    functions = [
        ("01", "稳定移载", "保持产品姿态并完成工位切换"),
        ("02", "姿态翻转", "按检测视角完成受控翻转"),
        ("03", "联锁确认", "到位、成像、放行信号闭环"),
    ]
    for index, (number, name, detail) in enumerate(functions):
        y = top + 1.60 + index * 0.58
        _add_text(slide, number, 8.16, y, 0.35, 0.25, size=7.5, color=COLORS["blue"], bold=True)
        _add_text(slide, name, 8.55, y, 1.05, 0.25, size=8.6, bold=True)
        _add_text(slide, detail, 9.65, y, 2.87, 0.28, size=7.7, color=COLORS["muted"])
    _add_key_value_rows(
        slide,
        "关键配置",
        [
            ("驱动形式", "伺服/气动组合"),
            ("定位方式", "专用治具 + 到位确认"),
            ("工作姿态", "俯视 / 仰视 / 侧视"),
            ("控制接口", "PLC状态交互"),
        ],
        8.16,
        top + 3.42,
        4.40,
        row_height=0.38,
    )
    _add_action_chain(
        slide,
        ["接收产品", "到位确认", "姿态切换", "信号确认", "释放移载"],
        5.92,
        left=0.58,
        width=7.24,
    )
    _add_footer(slide, 3, "模块配置为结构样本，正式选型以工程确认数据为准。")


def _build_module_detail(slide, images: dict[str, bytes], config: NavigationConfig) -> None:
    _add_page_title(slide, "3.2.2  搬运与翻转模块｜结构细节", "模块介绍 / MODULE DESCRIPTION", config)
    top = config.body_top
    _add_picture_contain(slide, images["handling"], 0.55, top, 7.02, 4.86, background=COLORS["light"])
    _add_caption(slide, "图 3.4  搬运与翻转模块整体结构", 0.55, top + 4.91, 7.02)

    _add_picture_contain(slide, images["loading"], 7.88, top, 4.82, 1.75, background=COLORS["light"])
    _add_caption(slide, "图 3.5  产品输入与定位衔接位", 7.88, top + 1.79, 4.82)
    _add_text(
        slide,
        "产品进入后完成缓存、基准定位和到位确认，为后续移载提供统一姿态。",
        7.94,
        top + 2.11,
        4.66,
        0.44,
        size=8.0,
        color=COLORS["muted"],
    )

    _add_picture_contain(slide, images["sorting"], 7.88, top + 2.76, 4.82, 1.70, background=COLORS["light"])
    _add_caption(slide, "图 3.6  产品输出与下料衔接位", 7.88, top + 4.50, 4.82)
    _add_text(
        slide,
        "检测完成后接收放行信号，并与分类下料模块完成节拍和异常状态交接。",
        7.94,
        top + 4.82,
        4.66,
        0.40,
        size=8.0,
        color=COLORS["muted"],
    )
    _add_footer(slide, 4, "细节页仅解释关键结构，不在此重复完整参数和检测效果。")


def _compact_text(value: str) -> str:
    return " ".join(str(value or "").replace("\n", " ").split())


def _build_detection_module(
    slide,
    images: dict[str, bytes],
    far: OpticalFarData,
    config: NavigationConfig,
) -> None:
    _add_page_title(
        slide,
        "3.3  检测模块介绍",
        "模块介绍 / INSPECTION MODULE",
        config,
    )
    top = config.body_top
    _add_picture_contain(
        slide,
        images["detection_module"],
        0.55,
        top,
        8.05,
        4.92,
        background=COLORS["light"],
    )
    marker_positions = (
        (1.66, top + 2.42),
        (4.30, top + 1.68),
        (6.47, top + 1.16),
    )
    for index, (left, marker_top) in enumerate(
        marker_positions[: len(far.stations)], start=1
    ):
        _add_number_marker(slide, f"{index:02d}", left, marker_top)
    _add_caption(
        slide,
        "图 3.7  检测模块结构与工位分布",
        0.55,
        top + 4.97,
        8.05,
    )

    _add_text(
        slide,
        "检测工位",
        8.92,
        top + 0.02,
        3.70,
        0.30,
        size=9.6,
        color=COLORS["blue"],
        bold=True,
    )
    _add_text(
        slide,
        f"本模块包含 {len(far.stations)} 个检测工位",
        8.92,
        top + 0.40,
        3.70,
        0.32,
        size=8.2,
        color=COLORS["muted"],
    )
    _add_line(slide, 8.92, top + 0.86, 3.70, color=COLORS["blue"], line_width=1.0)

    for index, station in enumerate(far.stations, start=1):
        view = _compact_text(station.view)
        station_name = f"{view}检测工位" if view else _compact_text(station.station)
        row_top = top + 1.18 + (index - 1) * 1.02
        _add_number_marker(slide, f"{index:02d}", 8.92, row_top)
        _add_text(
            slide,
            station_name or f"检测工位{index}",
            9.51,
            row_top - 0.01,
            3.05,
            0.32,
            size=10.5,
            bold=True,
            vertical=MSO_ANCHOR.MIDDLE,
        )
        _add_line(
            slide,
            8.92,
            row_top + 0.55,
            3.70,
            color=COLORS["line"],
            line_width=0.6,
        )

    _add_text(
        slide,
        "工位名称与模块图标注由工程师确认后写入正式方案。",
        8.92,
        top + 4.50,
        3.70,
        0.42,
        size=7.5,
        color=COLORS["muted"],
        align=PP_ALIGN.RIGHT,
    )
    _add_footer(slide, 5, "本页只介绍检测模块结构和工位名称；检测结果另章呈现。")


def _validate(
    output: Path,
    template_hash: str,
    navigation: NavigationConfig,
) -> None:
    if _sha256(TEMPLATE) != template_hash:
        raise RuntimeError("原始PPT模板在样本生成过程中发生变化")
    if output.read_bytes()[:4] != b"PK\x03\x04":
        raise RuntimeError("生成文件没有有效的Office ZIP签名")
    with zipfile.ZipFile(output) as package:
        names = set(package.namelist())
        required = {
            "[Content_Types].xml",
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
        }
        missing = sorted(required - names)
        if missing:
            raise RuntimeError(f"PPTX缺少必要成员：{missing}")
        if package.testzip() is not None:
            raise RuntimeError("PPTX ZIP成员损坏")
        slide_xml = [
            name for name in names if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        ]
        if len(slide_xml) < 5:
            raise RuntimeError(f"样本至少应包含5个Slide XML，实际为{len(slide_xml)}个")

    reopened = Presentation(output)
    if len(reopened.slides) != 5:
        raise RuntimeError(f"样本应为5页，实际为{len(reopened.slides)}页")
    expected_titles = (
        "设备总体方案",
        "功能模块分区",
        "搬运与翻转模块",
        "结构细节",
        "检测模块介绍",
    )
    picture_count = 0
    caption_count = 0
    rounded_shapes: list[str] = []
    out_of_bounds: list[str] = []
    font_names: set[str] = set()
    for slide_index, (slide, expected) in enumerate(
        zip(reopened.slides, expected_titles), start=1
    ):
        text = "\n".join(
            shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
        )
        if expected not in text:
            raise RuntimeError(f"样本缺少页面关键标题：{expected}")
        caption_count += sum(
            1
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.text.strip().startswith("图 3.")
        )
        picture_count += sum(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
        for shape in slide.shapes:
            if (
                shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                and shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
            ):
                rounded_shapes.append(f"第{slide_index}页:{shape.name}")
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > reopened.slide_width
                or shape.top + shape.height > reopened.slide_height
            ):
                out_of_bounds.append(f"第{slide_index}页:{shape.name}")
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            font_names.add(run.font.name)
    if picture_count < 7:
        raise RuntimeError(f"样本图片数量不足：{picture_count}")
    if caption_count < 7:
        raise RuntimeError(f"样本图题数量不足：{caption_count}")
    if rounded_shapes:
        raise RuntimeError(f"样本仍包含圆角矩形：{rounded_shapes}")
    if out_of_bounds:
        raise RuntimeError(f"样本存在越界对象：{out_of_bounds}")
    if font_names != {FONT_NAME}:
        raise RuntimeError(f"样本字体未统一：{sorted(font_names)}")

    detection_text = "\n".join(
        shape.text
        for shape in reopened.slides[4].shapes
        if getattr(shape, "has_text_frame", False)
    )
    for forbidden in ("OK样件", "NG", "相机", "镜头", "光源", "检测范围"):
        if forbidden in detection_text:
            raise RuntimeError(f"检测模块介绍页混入了检测效果或光学参数：{forbidden}")
    expected_stations = ("俯视检测工位", "仰视检测工位", "侧视检测工位")
    missing_stations = [
        station for station in expected_stations if station not in detection_text
    ]
    if missing_stations:
        raise RuntimeError(f"检测模块介绍页缺少工位名称：{missing_stations}")

    used_layouts = {slide.slide_layout.name for slide in reopened.slides}
    for layout in reopened.slide_layouts:
        if layout.name not in used_layouts:
            continue
        background = next(
            shape
            for shape in layout.shapes
            if shape.top < Inches(0.95)
            and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and shape.width > Inches(12.0)
        )
        actual_height = background.height / Inches(1)
        if abs(actual_height - navigation.height) > 0.01:
            raise RuntimeError(
                f"导航高度未生效：{layout.name}={actual_height:.3f} in"
            )
        actual_background = f"#{background.fill.fore_color.rgb}"
        if actual_background.upper() != "#FFFFFF":
            raise RuntimeError(
                f"导航栏基础背景未保持白色：{layout.name}={actual_background}"
            )
        effect_list = background.element.spPr.find(qn("a:effectLst"))
        outer_shadow = (
            effect_list.find(qn("a:outerShdw"))
            if effect_list is not None
            else None
        )
        if outer_shadow is not None:
            raise RuntimeError("导航背景灰色阴影未移除")
    for slide in reopened.slides:
        expected_font_size = PresentationStyle(
            navigation_height=navigation.height,
        ).resolved_navigation_font_size()
        navigation_text = [
            shape
            for shape in slide.shapes
            if shape.name.startswith("KY_NAV_TEXT_")
        ]
        if tuple(shape.text.strip() for shape in navigation_text) != NAVIGATION_ITEMS:
            raise RuntimeError("导航栏目名称或顺序不正确")
        active_text = navigation_text[2]
        active_background = next(
            shape for shape in slide.shapes if shape.name == "KY_NAV_ACTIVE_BACKGROUND"
        )
        active_bar = next(
            shape for shape in slide.shapes if shape.name == "KY_NAV_ACTIVE"
        )
        baseline = next(
            shape for shape in slide.shapes if shape.name == "KY_NAV_BASELINE"
        )
        if f"#{active_background.fill.fore_color.rgb}".upper() != navigation.background.upper():
            raise RuntimeError("当前导航栏目背景色未生效")
        if f"#{baseline.fill.fore_color.rgb}".upper() != "#D3D9DE":
            raise RuntimeError("导航灰色底线颜色不正确")
        if (
            active_background.left != active_text.left
            or active_background.width != active_text.width
            or active_bar.left != active_text.left
            or active_bar.width != active_text.width
            or baseline.left != 0
            or abs(baseline.left + baseline.width - reopened.slide_width) > 1
            or baseline.top != active_bar.top
            or baseline.height != active_bar.height
        ):
            raise RuntimeError("导航底线或当前栏目红线位置不正确")
        baseline_effect = baseline.element.spPr.find(qn("a:effectLst"))
        if (
            baseline_effect is not None
            and baseline_effect.find(qn("a:outerShdw")) is not None
        ):
            raise RuntimeError("导航灰色底线不应包含阴影")
        if any(
            shape.name.startswith("KY_NAV_INACTIVE_")
            for shape in slide.shapes
        ):
            raise RuntimeError("不应生成逐栏白色条带")
        if navigation_text[0].left != 0:
            raise RuntimeError("第一导航栏必须从幻灯片左边界开始")
        for previous, current in zip(navigation_text, navigation_text[1:]):
            if abs(previous.left + previous.width - current.left) > 1:
                raise RuntimeError("相邻导航栏目之间存在留白或重叠")
        for shape in navigation_text:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.font.bold:
                        raise RuntimeError(f"导航文字未加粗：{shape.text}")
                    if run.font.size is None or abs(
                        run.font.size.pt - expected_font_size
                    ) > 0.01:
                        raise RuntimeError(f"导航文字字号不正确：{shape.text}")


def build(
    output: Path,
    *,
    overwrite: bool = False,
    navigation: NavigationConfig | None = None,
) -> Path:
    navigation = navigation or NavigationConfig()
    navigation.validate()
    if output.exists() and not overwrite:
        raise FileExistsError(f"输出文件已存在：{output}")
    output.parent.mkdir(parents=True, exist_ok=True)

    template_hash = _sha256(TEMPLATE)
    presentation = Presentation(TEMPLATE)
    far = parse_optical_far(FAR_SOURCE)
    if not far.stations:
        raise RuntimeError("光学FAR缺少工位数据，无法制作检测模块介绍页")
    images = {
        "overview": _picture_blob(presentation, 5, 3),
        "module_map": _picture_blob(presentation, 6, 20),
        "loading": _picture_blob(presentation, 7, 7),
        "handling": _picture_blob(presentation, 8, 6),
        "sorting": _picture_blob(presentation, 9, 2),
        "detection_module": _picture_blob(presentation, 8, 6),
    }

    equipment_layout = _find_layout(presentation, "2_标题幻灯片")
    module_layout = _find_layout(presentation, "3_标题幻灯片")
    slides = [
        presentation.slides.add_slide(equipment_layout),
        presentation.slides.add_slide(equipment_layout),
        presentation.slides.add_slide(module_layout),
        presentation.slides.add_slide(module_layout),
        presentation.slides.add_slide(module_layout),
    ]
    _remove_other_slides(presentation, slides)
    for slide in slides:
        _clear_slide(slide)
    _configure_navigation(presentation, navigation, [2] * len(slides))

    _build_overview(slides[0], images, navigation)
    _build_module_map(slides[1], images, navigation)
    _build_module_intro(slides[2], images, navigation)
    _build_module_detail(slides[3], images, navigation)
    _build_detection_module(slides[4], images, far, navigation)

    presentation.save(output)
    _validate(output, template_hash, navigation)
    return output.resolve()


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--overwrite", action="store_true")
    parser.add_argument(
        "--nav-height",
        type=float,
        default=NavigationConfig().height,
        help="顶部导航栏高度（英寸），允许范围0.42～0.72，默认0.52",
    )
    parser.add_argument(
        "--nav-background",
        default=NavigationConfig().background,
        help="当前导航栏目背景色，使用#RRGGBB格式，默认#FFFFFF",
    )
    args = parser.parse_args()
    navigation = NavigationConfig(
        height=args.nav_height,
        background=args.nav_background.upper(),
    )
    print(
        build(
            args.output.resolve(),
            overwrite=args.overwrite,
            navigation=navigation,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
