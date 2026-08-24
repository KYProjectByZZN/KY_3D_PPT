"""Build a reproducible equipment-flow/module validation project and PPTX."""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from ppt_generator.module_service import ensure_project_modules
from ppt_generator.project import DeviceModule, FlowNode, PptProject, save_project
from ppt_generator.scheme_service import materialize_equipment_scheme
from ppt_generator.template_renderer import load_manifest, render_project


TEMPLATE = PROJECT_ROOT / "templates" / "冲压筒形壳体检测方案NAT6704_v2.pptx"
MANIFEST = PROJECT_ROOT / "templates" / "NAT6704_v2.template.json"
DEFAULT_DATA = PROJECT_ROOT / "examples" / "NAT6704_v2_test_data.json"
OUTPUT_PPTX = PROJECT_ROOT / "output" / "设备方案流程模块_v0.7.0_验收版.pptx"
OUTPUT_PROJECT = PROJECT_ROOT / "output" / "设备方案流程模块_v0.7.0_验收版.pptproj.json"
ASSET_ROOT = PROJECT_ROOT / "output" / "scheme_validation_assets_v070_acceptance"


def _extract_picture(
    presentation: Presentation,
    slide_number: int,
    shape_id: int,
    stem: str,
) -> Path:
    shape = next(
        item
        for item in presentation.slides[slide_number - 1].shapes
        if item.shape_id == shape_id
    )
    image = shape.image
    path = ASSET_ROOT / f"{stem}.{image.ext}"
    path.write_bytes(image.blob)
    return path.resolve()


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--overwrite", action="store_true")
    args = parser.parse_args()
    for path in (OUTPUT_PPTX, OUTPUT_PROJECT):
        if path.exists() and not args.overwrite:
            raise FileExistsError(f"验证产物已存在：{path}")

    ASSET_ROOT.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(TEMPLATE)
    overview_image = _extract_picture(presentation, 5, 3, "equipment_overview")
    module_images = [
        _extract_picture(presentation, 6, 20, "feed_and_layout"),
        _extract_picture(presentation, 7, 7, "handling"),
        _extract_picture(presentation, 8, 6, "vision_and_transfer"),
        _extract_picture(presentation, 9, 2, "sorting_and_output"),
    ]

    values = json.loads(DEFAULT_DATA.read_text(encoding="utf-8"))
    project = PptProject(
        project_name="设备方案流程与模块结构验证",
        template_path=str(TEMPLATE),
        manifest_path=str(MANIFEST),
        output_path=str(OUTPUT_PPTX),
        values=values,
    )
    manifest = load_manifest(MANIFEST)
    ensure_project_modules(project, manifest)
    project.equipment_scheme.initialized = True
    project.equipment_scheme.overview_image = str(overview_image)
    project.equipment_scheme.overview_description = (
        "本验证设备由上料定位、搬运翻转、视觉检测和分类下料模块组成，"
        "各模块按照确认的检测流程协同运行。"
    )

    devices = [
        DeviceModule(
            name="上料定位模块",
            module_type="上料",
            station="上料位",
            function="完成产品上料、缓存和检测前定位。",
            action="接收产品后定位，并向检测系统输出到位信号。",
            image_path=str(module_images[0]),
            page_template_key="equipment_module_page_6",
        ),
        DeviceModule(
            name="搬运翻转模块",
            module_type="翻转",
            station="工位1～工位3",
            function="完成产品在多个检测工位间的搬运与姿态切换。",
            action="依次完成移载、翻转和检测位置确认。",
            image_path=str(module_images[1]),
            page_template_key="equipment_module_page_7",
        ),
        DeviceModule(
            name="视觉检测模块",
            module_type="视觉检测",
            station="工位1～工位3",
            function="完成俯视、仰视和侧视成像检测。",
            action="收到触发信号后采图、分析并输出检测结果。",
            image_path=str(module_images[2]),
            page_template_key="equipment_module_page_8",
        ),
        DeviceModule(
            name="分类下料模块",
            module_type="下料",
            station="下料位",
            function="按照检测结果完成OK和NG产品分类。",
            action="接收结果并把产品送入对应下料位置。",
            image_path=str(module_images[3]),
            page_template_key="equipment_module_page_9",
        ),
    ]
    project.equipment_scheme.equipment_modules = devices
    project.equipment_scheme.flow_nodes = [
        FlowNode("人工上料", "上料", "上料位", "放入待检产品", devices[0].id),
        FlowNode("产品定位", "定位", "上料位", "定位并确认到位", devices[0].id),
        FlowNode("工位1俯视检测", "检测", "工位1", "完成俯视成像", devices[2].id),
        FlowNode("产品翻转", "翻转", "翻转位", "切换产品姿态", devices[1].id),
        FlowNode("工位2仰视检测", "检测", "工位2", "完成仰视成像", devices[2].id),
        FlowNode("工位3侧视检测", "检测", "工位3", "完成侧视成像", devices[2].id),
        FlowNode("结果判定", "分拣", "判定节点", "汇总检测结果", devices[3].id),
        FlowNode("分类下料", "下料", "下料位", "OK/NG分类输出", devices[3].id),
    ]

    result = materialize_equipment_scheme(project, manifest)
    save_project(project, OUTPUT_PROJECT)
    render_project(project, overwrite=args.overwrite)
    print(
        f"{OUTPUT_PPTX.resolve()} | flow={result.flow_pages} | "
        f"equipment={result.equipment_pages}"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
