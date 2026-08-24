"""Build an editable reference pack for industrial equipment PPT styling.

The pack uses selected pages from official vendor PDFs as small, attributed
reference screenshots. It does not modify the production template or renderer.
"""

from __future__ import annotations

import argparse
import zipfile
from io import BytesIO
from pathlib import Path

import fitz
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
ASSET_ROOT = PROJECT_ROOT / "output" / "industrial_reference_assets"
RENDER_ROOT = ASSET_ROOT / "selected_pages"
DEFAULT_OUTPUT = PROJECT_ROOT / "output" / "设备介绍_工业版式参考与工位页方案_v1.pptx"

FONT_NAME = "Microsoft YaHei"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1F, 0x25, 0x2B)
MUTED = RGBColor(0x62, 0x6B, 0x73)
LINE = RGBColor(0xC7, 0xCD, 0xD2)
LIGHT = RGBColor(0xF2, 0xF4, 0xF5)
LIGHT_BLUE = RGBColor(0xE8, 0xEF, 0xF4)
BLUE = RGBColor(0x2B, 0x58, 0x75)
RED = RGBColor(0xC9, 0x00, 0x00)
GREEN = RGBColor(0x2D, 0x6A, 0x4F)

SOURCES = {
    "beckhoff": {
        "pdf": ASSET_ROOT / "beckhoff_xplanar.pdf",
        "title": "Beckhoff｜XPlanar: Levitating, contactless, intelligent",
        "url": "https://download.beckhoff.com/download/document/catalog/Beckhoff_XPlanar_e.pdf",
    },
    "omron_pharma": {
        "pdf": ASSET_ROOT / "omron_fh_flyer.pdf",
        "title": "OMRON｜Inspection Solutions for Pharma, FH Series",
        "url": "https://assets.omron.eu/downloads/latest/brochure/en/y269_fh_series_vision_system_flyer_en.pdf?v=1",
    },
    "omron_system": {
        "pdf": ASSET_ROOT / "omron_fh_inspection.pdf",
        "title": "OMRON｜Inspection systems, FH series",
        "url": "https://assets.omron.eu/downloads/latest/brochure/en/xpectia_fh_brochure_en.pdf",
    },
    "sick": {
        "pdf": ASSET_ROOT / "sick_nova_inspection.pdf",
        "title": "SICK｜Nova Intelligent Inspection",
        "url": "https://www.sick.com/media/docs/2/82/582/flyer_sick_nova_intelligent_inspection_2d_machine_vision_en_im0101582.pdf",
    },
    "omron_pcb": {
        "pdf": ASSET_ROOT / "omron_pcb_inspection.pdf",
        "title": "OMRON｜PCB Inspection System",
        "url": "https://assets.omron.eu/downloads/latest/brochure/en/q344_vt-s1080_pcb_inspection_system_brochure_en.pdf?v=1",
    },
    "keyence": {
        "pdf": ASSET_ROOT / "keyence_corporate_profile.pdf",
        "title": "KEYENCE｜Corporate Profile / Machine Vision Systems",
        "url": "https://www.keyence.com/pdf/corporate_profile.pdf",
    },
}

SELECTED_PAGES = {
    "beckhoff_components": ("beckhoff", 4),
    "beckhoff_applications": ("beckhoff", 8),
    "omron_line_a": ("omron_pharma", 1),
    "omron_line_b": ("omron_pharma", 2),
    "omron_station_scope": ("omron_pharma", 3),
    "omron_station_specs": ("omron_pharma", 4),
    "omron_system_config": ("omron_system", 12),
    "sick_examples": ("sick", 1),
    "omron_result": ("omron_pcb", 3),
    "keyence_vision": ("keyence", 10),
}


def _apply_font(run, *, size: float, color: RGBColor, bold: bool = False) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        existing = r_pr.find(qn(tag))
        if existing is None:
            existing = OxmlElement(tag)
            r_pr.append(existing)
        existing.set("typeface", FONT_NAME)


def _add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: float = 11.0,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.04,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = 1.05
    run = paragraph.add_run()
    run.text = text
    _apply_font(run, size=size, color=color, bold=bold)
    return box


def _add_paragraphs(
    slide,
    rows: list[tuple[str, bool]],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: float = 10.5,
    color: RGBColor = INK,
    gap: float = 4.0,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, (text, bold) in enumerate(rows):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.08
        run = paragraph.add_run()
        run.text = text
        _apply_font(run, size=size, color=color, bold=bold)
    return box


def _add_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = LINE,
    line_width: float = 0.8,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def _add_rule(slide, left: float, top: float, width: float, color: RGBColor = LINE):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(left),
        Inches(top),
        Inches(left + width),
        Inches(top),
    )
    line.line.color.rgb = color
    line.line.width = Pt(0.8)
    return line


def _add_title(slide, section: str, title: str, subtitle: str = "") -> None:
    _add_text(slide, section, 0.48, 0.26, 0.78, 0.32, size=9.0, color=RED, bold=True)
    _add_text(slide, title, 1.18, 0.20, 10.7, 0.48, size=22.0, color=INK, bold=True)
    if subtitle:
        _add_text(slide, subtitle, 1.20, 0.68, 10.9, 0.30, size=8.7, color=MUTED)
    _add_rule(slide, 0.48, 1.03, 12.37, RED)


def _add_footer(slide, page_number: int, source: str | None = None) -> None:
    if source:
        _add_text(slide, source, 0.50, 7.18, 11.85, 0.20, size=6.8, color=MUTED)
    _add_text(
        slide,
        f"{page_number:02d}",
        12.76,
        7.24,
        0.52,
        0.24,
        size=7.8,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def _render_selected_pages() -> dict[str, Path]:
    RENDER_ROOT.mkdir(parents=True, exist_ok=True)
    rendered: dict[str, Path] = {}
    documents: dict[str, fitz.Document] = {}
    try:
        for output_key, (source_key, page_index) in SELECTED_PAGES.items():
            source = SOURCES[source_key]
            pdf_path = Path(source["pdf"])
            if not pdf_path.exists():
                raise FileNotFoundError(f"Missing reference PDF: {pdf_path}")
            document = documents.setdefault(source_key, fitz.open(pdf_path))
            if page_index >= document.page_count:
                raise ValueError(f"Page {page_index + 1} outside {pdf_path.name}")
            output = RENDER_ROOT / f"{output_key}.png"
            if not output.exists() or output.stat().st_mtime < pdf_path.stat().st_mtime:
                pixmap = document[page_index].get_pixmap(
                    matrix=fitz.Matrix(1.8, 1.8), alpha=False
                )
                pixmap.save(output)
            rendered[output_key] = output
    finally:
        for document in documents.values():
            document.close()
    return rendered


def _add_picture_contain(
    slide,
    image_path: Path,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    border: bool = True,
):
    if border:
        _add_rect(slide, left, top, width, height, fill=WHITE, line=LINE, line_width=0.7)
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    padding = 0.05
    available_width = width - padding * 2
    available_height = height - padding * 2
    scale = min(available_width / image_width, available_height / image_height)
    display_width = image_width * scale
    display_height = image_height * scale
    return slide.shapes.add_picture(
        str(image_path),
        Inches(left + (width - display_width) / 2),
        Inches(top + (height - display_height) / 2),
        Inches(display_width),
        Inches(display_height),
    )


def _add_caption(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    *,
    size: float = 8.5,
):
    _add_text(
        slide,
        text,
        left,
        top,
        width,
        0.25,
        size=size,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def _add_source_box(slide, source_key: str, left: float, top: float, width: float):
    source = SOURCES[source_key]
    _add_text(slide, str(source["title"]), left, top, width, 0.24, size=7.8, color=INK, bold=True)
    box = _add_text(slide, str(source["url"]), left, top + 0.23, width, 0.36, size=6.4, color=BLUE)
    for paragraph in box.text_frame.paragraphs:
        for run in paragraph.runs:
            run.hyperlink.address = str(source["url"])


def _add_number_row(
    slide,
    number: str,
    title: str,
    description: str,
    top: float,
    *,
    accent: RGBColor = BLUE,
):
    _add_rect(slide, 0.70, top, 0.54, 0.54, fill=accent, line=accent)
    _add_text(
        slide,
        number,
        0.70,
        top,
        0.54,
        0.54,
        size=10.0,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )
    _add_text(slide, title, 1.48, top - 0.01, 2.65, 0.27, size=11.5, bold=True)
    _add_text(slide, description, 4.05, top - 0.01, 8.05, 0.44, size=9.3, color=MUTED)
    _add_rule(slide, 1.48, top + 0.60, 10.70)


def _slide_cover(prs: Presentation, images: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, 13.333, 7.5, fill=WHITE, line=WHITE, line_width=0)
    _add_rect(slide, 0, 0, 0.18, 7.5, fill=RED, line=RED, line_width=0)
    _add_text(slide, "工业设备技术方案 PPT", 0.72, 0.72, 6.0, 0.36, size=11.0, color=RED, bold=True)
    _add_text(
        slide,
        "设备介绍版式参考\n与检测工位页方案",
        0.68,
        1.26,
        6.2,
        1.18,
        size=28.0,
        color=INK,
        bold=True,
    )
    _add_text(
        slide,
        "结论：取消圆角卡片化表达，改为工程图主导、直角分区、线性标注和规范图题。",
        0.72,
        2.66,
        5.65,
        0.68,
        size=12.0,
        color=BLUE,
        bold=True,
    )
    _add_paragraphs(
        slide,
        [
            ("字体：中文、英文、数字统一使用 Microsoft YaHei", False),
            ("图片：不做圆角蒙版，不做胶囊标签，不加重阴影", False),
            ("图题：每张图下方显示“图 3.1 设备总体图”", False),
            ("新增：检测模块总览 + 每工位一页检测工位详页", False),
        ],
        0.72,
        3.55,
        5.85,
        1.70,
        size=10.3,
        color=INK,
        gap=8,
    )
    _add_picture_contain(slide, images["beckhoff_components"], 6.95, 0.65, 5.80, 2.72)
    _add_picture_contain(slide, images["omron_line_a"], 6.95, 3.55, 2.78, 2.88)
    _add_picture_contain(slide, images["omron_station_scope"], 9.95, 3.55, 2.80, 2.88)
    _add_caption(slide, "官方工业资料截图缩略图，仅用于版式研究", 6.95, 6.52, 5.80)
    _add_text(slide, "待确认稿｜2026-08-24", 0.72, 6.92, 4.4, 0.26, size=8.5, color=MUTED)
    _add_footer(slide, 1)


def _slide_rules(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "01", "先修正五条基础规范", "这些规则应先于具体页面样式执行")
    rows = [
        ("01", "取消圆角容器", "设备图、CAD图、参数区统一使用直角边界；必要时只保留0.5～0.8 pt细线。"),
        ("02", "统一一种字体", "中文、英文、数字全部使用 Microsoft YaHei；不再混用微软雅黑、Arial、Aptos。"),
        ("03", "建立规范图题", "主图下方使用“图 3.1 设备总体图”；局部图按图号继续递增。"),
        ("04", "减少装饰性标签", "取消胶囊标签和大面积浅色卡片，编号只服务定位和结构关系。"),
        ("05", "明确页面职责", "结构页讲设备与工位怎么组成；检测效果页只展示OK/NG证据，不重复规格。"),
    ]
    for index, row in enumerate(rows):
        _add_number_row(slide, *row, 1.34 + index * 1.02, accent=BLUE if index < 4 else RED)
    _add_text(
        slide,
        "建议图题字号 9～10 pt，页面标题 22～24 pt，正文 10.5～12 pt；最小正文不低于10 pt。",
        0.74,
        6.60,
        11.7,
        0.38,
        size=9.0,
        color=MUTED,
    )
    _add_footer(slide, 2)


def _slide_beckhoff(prs: Presentation, images: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "02", "参考一：整机/模块页采用大图 + 直线标注", "Beckhoff XPlanar 官方资料的工程表达方式")
    _add_picture_contain(slide, images["beckhoff_components"], 0.62, 1.30, 8.06, 4.92)
    _add_caption(slide, "参考图 2.1  系统构成与直线标注", 0.62, 6.26, 8.06)
    _add_text(slide, "可直接借鉴", 9.03, 1.34, 3.30, 0.32, size=12.0, color=RED, bold=True)
    _add_paragraphs(
        slide,
        [
            ("主图占页面主体，信息层级先图后文。", False),
            ("标注线直接指向真实部件，不把说明装进圆角卡片。", False),
            ("部件图、系统图保持完整比例，背景只做浅灰承托。", False),
            ("参数区是小型表格/列表，不使用大号数字卡片。", False),
        ],
        9.00,
        1.82,
        3.55,
        2.12,
        size=10.2,
        gap=7,
    )
    _add_rect(slide, 9.03, 4.24, 3.38, 1.45, fill=LIGHT, line=LINE)
    _add_text(slide, "用于本项目", 9.24, 4.45, 1.45, 0.27, size=9.5, color=BLUE, bold=True)
    _add_text(
        slide,
        "设备总体图\n功能模块分区图\n单模块结构介绍页",
        9.24,
        4.80,
        2.80,
        0.76,
        size=10.0,
        color=INK,
    )
    _add_source_box(slide, "beckhoff", 9.02, 5.93, 3.55)
    _add_footer(slide, 3, "官方来源见右下角链接；截图仅用于内部版式研究。")


def _slide_applications(prs: Presentation, images: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "03", "参考二：同类设备图用平铺，不做圆角画廊", "产品变体与应用图使用统一基线、同尺度和图下说明")
    _add_picture_contain(slide, images["beckhoff_applications"], 0.62, 1.30, 7.22, 4.72)
    _add_caption(slide, "参考图 3.1  三类应用场景及线性图注", 0.62, 6.06, 7.22)
    _add_picture_contain(slide, images["omron_system_config"], 8.17, 1.30, 4.52, 3.48)
    _add_caption(slide, "参考图 3.2  视觉系统构成与连接关系", 8.17, 4.82, 4.52)
    _add_text(
        slide,
        "判断：多图页可以保留，但必须同宽、同高、同基线；每张图都有明确图题。图片本身不做圆角，图之间靠留白和细线区分。",
        8.20,
        5.28,
        4.42,
        0.82,
        size=9.7,
        color=INK,
    )
    _add_source_box(slide, "beckhoff", 0.65, 6.48, 6.80)
    _add_source_box(slide, "omron_system", 7.80, 6.48, 4.82)
    _add_footer(slide, 4)


def _slide_station_overview(prs: Presentation, images: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "04", "参考三：先给检测模块总览，再进入每个工位", "OMRON 官方资料按“产线概览 → 工位详页”组织内容")
    _add_picture_contain(slide, images["omron_line_a"], 0.62, 1.30, 5.95, 4.55)
    _add_picture_contain(slide, images["omron_line_b"], 6.80, 1.30, 5.90, 4.55)
    _add_caption(slide, "参考图 4.1  产线前半段工位总览", 0.62, 5.90, 5.95)
    _add_caption(slide, "参考图 4.2  产线后半段工位总览", 6.80, 5.90, 5.90)
    _add_text(
        slide,
        "对本项目的转换：检测模块总览只负责工位编号、视角、产品姿态和工位关系；当仅1个工位时可省略总览，直接进入工位详页。",
        0.72,
        6.34,
        8.65,
        0.54,
        size=9.8,
        color=INK,
        bold=True,
    )
    _add_source_box(slide, "omron_pharma", 9.55, 6.26, 3.02)
    _add_footer(slide, 5)


def _slide_station_detail(prs: Presentation, images: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "05", "参考四：单工位内容按“检测范围 + 光学结构 + 参数”组织", "官方资料分成两页；本项目压缩成一页，避免方案过长")
    _add_picture_contain(slide, images["omron_station_scope"], 0.62, 1.27, 5.08, 4.66)
    _add_picture_contain(slide, images["omron_station_specs"], 5.94, 1.27, 5.08, 4.66)
    _add_caption(slide, "参考图 5.1  工位检测范围与机构示意", 0.62, 5.97, 5.08)
    _add_caption(slide, "参考图 5.2  工位特点、参数与光学结构", 5.94, 5.97, 5.08)
    _add_rect(slide, 11.24, 1.28, 1.46, 4.65, fill=LIGHT_BLUE, line=BLUE)
    _add_text(slide, "本项目\n单页字段", 11.43, 1.53, 1.06, 0.68, size=10.0, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    _add_text(
        slide,
        "工位/视角\n检测项\n产品姿态\n相机\n镜头\n光源\n触发方式\n节拍\n输入/输出",
        11.42,
        2.42,
        1.08,
        2.93,
        size=8.6,
        color=INK,
        align=PP_ALIGN.CENTER,
    )
    _add_source_box(slide, "omron_pharma", 0.65, 6.38, 7.85)
    _add_text(slide, "原则：一个工位一页；字段缺失则提示补充，不用装饰性内容填空。", 8.60, 6.46, 3.98, 0.34, size=9.0, color=RED, bold=True)
    _add_footer(slide, 6)


def _slide_boundaries(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "06", "不会冲突，但要把四类页面的职责切开", "检测工位详页替代通用“视觉检测模块页”，而不是重复再加一份")
    headers = ["页面类型", "只回答的问题", "允许内容", "禁止重复"]
    widths = [2.05, 3.18, 3.70, 3.18]
    x = 0.62
    for index, (header, width) in enumerate(zip(headers, widths)):
        _add_rect(slide, x, 1.38, width, 0.56, fill=BLUE if index == 0 else LIGHT_BLUE, line=LINE)
        _add_text(
            slide,
            header,
            x + 0.08,
            1.49,
            width - 0.16,
            0.28,
            size=9.2,
            color=WHITE if index == 0 else INK,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        x += width
    rows = [
        ("设备功能模块页", "机构做什么、如何动作", "机械结构、动作链、接口", "相机型号、缺陷样图"),
        ("检测模块总览", "有几个工位、怎么布置", "工位编号、视角、产品姿态", "逐工位参数、缺陷大图"),
        ("检测工位详页", "该工位怎样完成检测", "光学结构、检测项、触发、节拍", "大量OK/NG结果图"),
        ("检测效果页", "实际能看到什么缺陷", "OK/NG样图、缺陷说明", "机械动作、重复光学清单"),
        ("视觉系统总页", "共用平台如何组成", "控制器、HMI、算法架构、网络", "每工位重复型号与参数"),
    ]
    y = 1.94
    for row_index, row in enumerate(rows):
        x = 0.62
        fill = WHITE if row_index % 2 == 0 else LIGHT
        for col_index, (value, width) in enumerate(zip(row, widths)):
            _add_rect(slide, x, y, width, 0.80, fill=fill, line=LINE)
            _add_text(
                slide,
                value,
                x + 0.10,
                y + 0.12,
                width - 0.20,
                0.55,
                size=8.8,
                color=BLUE if col_index == 0 else INK,
                bold=col_index == 0,
                align=PP_ALIGN.CENTER if col_index == 0 else PP_ALIGN.LEFT,
                valign=MSO_ANCHOR.MIDDLE,
            )
            x += width
        y += 0.80
    _add_rect(slide, 0.62, 6.20, 12.11, 0.52, fill=LIGHT_BLUE, line=BLUE)
    _add_text(
        slide,
        "NAT6801示例：3个检测工位 → 1页检测模块总览 + 3页工位详页；现有6页检测效果继续保留。",
        0.84,
        6.32,
        11.67,
        0.28,
        size=9.7,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_footer(slide, 7)


def _slide_result_reference(prs: Presentation, images: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "07", "检测效果页与工位详页必须分开", "厂商资料也把应用证据、检测工具和系统配置分成不同页面")
    _add_picture_contain(slide, images["sick_examples"], 0.62, 1.30, 5.90, 4.55)
    _add_picture_contain(slide, images["omron_result"], 6.80, 1.30, 5.90, 4.55)
    _add_caption(slide, "参考图 7.1  多种检测应用实例", 0.62, 5.90, 5.90)
    _add_caption(slide, "参考图 7.2  检测原理与结果证据", 6.80, 5.90, 5.90)
    _add_text(slide, "保留到检测效果页", 0.72, 6.34, 2.00, 0.27, size=9.3, color=GREEN, bold=True)
    _add_text(slide, "OK/NG样图、缺陷区域、判定说明", 2.60, 6.34, 3.80, 0.27, size=9.3, color=INK)
    _add_text(slide, "移出检测效果页", 6.92, 6.34, 2.00, 0.27, size=9.3, color=RED, bold=True)
    _add_text(slide, "完整光学清单、机械动作、通用平台说明", 8.80, 6.34, 3.85, 0.27, size=9.3, color=INK)
    _add_source_box(slide, "sick", 0.67, 6.72, 5.72)
    _add_source_box(slide, "omron_pcb", 6.73, 6.72, 5.88)
    _add_footer(slide, 8)


def _slide_wireframes(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, "08", "推荐采用的最终页面骨架", "方案A：Beckhoff式工程主图 + OMRON式单工位结构（推荐）")

    _add_text(slide, "A｜设备/模块页", 0.64, 1.23, 3.0, 0.31, size=11.5, color=RED, bold=True)
    _add_rect(slide, 0.62, 1.65, 5.94, 4.86, fill=WHITE, line=INK, line_width=1.0)
    _add_text(slide, "3.1 设备总体方案", 0.88, 1.88, 4.65, 0.36, size=14.0, bold=True)
    _add_rule(slide, 0.88, 2.33, 5.18, RED)
    _add_rect(slide, 0.92, 2.56, 3.54, 2.75, fill=LIGHT, line=LINE)
    _add_text(slide, "整机图 / CAD等轴测图\n完整显示，不裁切，不圆角", 1.38, 3.52, 2.62, 0.78, size=11.0, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    _add_caption(slide, "图 3.1  设备总体图", 0.92, 5.36, 3.54, size=9.0)
    _add_text(slide, "设备名称", 4.72, 2.60, 1.28, 0.28, size=9.0, color=MUTED)
    _add_text(slide, "冲压筒形壳体视觉检测设备", 4.72, 2.91, 1.58, 0.64, size=11.0, bold=True)
    _add_rule(slide, 4.72, 3.68, 1.48)
    for idx, (name, value) in enumerate((("检测工位", "3"), ("生产节拍", "40–50 pcs/min"), ("上料方式", "人工/输送线"))):
        y = 3.90 + idx * 0.57
        _add_text(slide, name, 4.72, y, 0.72, 0.26, size=8.0, color=MUTED)
        _add_text(slide, value, 5.35, y, 0.88, 0.28, size=8.5, color=INK, bold=True, align=PP_ALIGN.RIGHT)

    _add_text(slide, "B｜单工位检测详页", 6.82, 1.23, 3.4, 0.31, size=11.5, color=RED, bold=True)
    _add_rect(slide, 6.80, 1.65, 5.94, 4.86, fill=WHITE, line=INK, line_width=1.0)
    _add_text(slide, "4.2 工位1：俯视检测", 7.06, 1.88, 4.65, 0.36, size=14.0, bold=True)
    _add_rule(slide, 7.06, 2.33, 5.18, RED)
    _add_rect(slide, 7.10, 2.56, 3.18, 2.57, fill=LIGHT, line=LINE)
    _add_text(slide, "工位机构图\n相机视角 + 产品姿态 + 照明位置", 7.48, 3.38, 2.42, 0.78, size=10.0, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    _add_caption(slide, "图 4.1  工位1俯视检测模块", 7.10, 5.18, 3.18, size=9.0)
    _add_text(slide, "检测范围", 10.58, 2.56, 0.90, 0.25, size=8.2, color=BLUE, bold=True)
    _add_text(slide, "指纹、镀层不良、生锈、磨伤", 10.58, 2.87, 1.86, 0.56, size=8.8, color=INK)
    _add_rule(slide, 10.58, 3.54, 1.75)
    fields = (("相机", "1200万像素"), ("镜头", "FA镜头"), ("光源", "同轴光"), ("触发", "到位信号"), ("节拍", "≤1.5 s"))
    for idx, (name, value) in enumerate(fields):
        y = 3.72 + idx * 0.38
        _add_text(slide, name, 10.58, y, 0.54, 0.22, size=7.8, color=MUTED)
        _add_text(slide, value, 11.13, y, 1.19, 0.22, size=7.8, color=INK, bold=True, align=PP_ALIGN.RIGHT)

    _add_rect(slide, 0.62, 6.76, 12.12, 0.42, fill=LIGHT_BLUE, line=BLUE)
    _add_text(
        slide,
        "建议确认：采用方案A；现有A～D样本按此骨架重做，并新增1张工位详页样本。",
        0.90,
        6.85,
        11.55,
        0.24,
        size=9.8,
        color=BLUE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_footer(slide, 9)


def build(output: Path, *, overwrite: bool = False) -> None:
    if output.exists() and not overwrite:
        raise FileExistsError(f"Output already exists: {output}")
    output.parent.mkdir(parents=True, exist_ok=True)
    images = _render_selected_pages()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    _slide_cover(prs, images)
    _slide_rules(prs)
    _slide_beckhoff(prs, images)
    _slide_applications(prs, images)
    _slide_station_overview(prs, images)
    _slide_station_detail(prs, images)
    _slide_boundaries(prs)
    _slide_result_reference(prs, images)
    _slide_wireframes(prs)
    prs.save(output)

    if output.read_bytes()[:2] != b"PK":
        raise RuntimeError("Generated file is not an Office ZIP package")
    with zipfile.ZipFile(output) as archive:
        required = {"[Content_Types].xml", "ppt/presentation.xml"}
        missing = required.difference(archive.namelist())
        if missing:
            raise RuntimeError(f"Generated PPTX missing entries: {sorted(missing)}")
    reopened = Presentation(output)
    if len(reopened.slides) != 9:
        raise RuntimeError(f"Expected 9 slides, found {len(reopened.slides)}")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--overwrite", action="store_true")
    args = parser.parse_args()
    build(args.output, overwrite=args.overwrite)
    print(args.output)


if __name__ == "__main__":
    main()
