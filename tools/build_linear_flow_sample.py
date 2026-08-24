"""Build a one-slide editable industrial linear-flow PPT sample."""

from __future__ import annotations

import argparse
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = PROJECT_ROOT / "templates" / "冲压筒形壳体检测方案NAT6704_v2.pptx"
DEFAULT_OUTPUT = PROJECT_ROOT / "output" / "设备单主线流程图_样本_v1.pptx"

FONT_NAME = "Microsoft YaHei"
COLORS = {
    "background": RGBColor(0xF5, 0xF7, 0xF9),
    "card": RGBColor(0xFF, 0xFF, 0xFF),
    "card_blue": RGBColor(0xF3, 0xF8, 0xFC),
    "card_edge": RGBColor(0xD5, 0xDE, 0xE6),
    "primary": RGBColor(0x1F, 0x4E, 0x79),
    "vision": RGBColor(0x2F, 0x75, 0xB5),
    "transfer": RGBColor(0x60, 0x70, 0x7D),
    "connector": RGBColor(0x8A, 0x99, 0xA8),
    "text": RGBColor(0x1F, 0x29, 0x33),
    "muted": RGBColor(0x66, 0x73, 0x7F),
    "soft_line": RGBColor(0xE4, 0xE9, 0xEE),
    "ok_fill": RGBColor(0xE7, 0xF4, 0xEA),
    "ok_text": RGBColor(0x1B, 0x5E, 0x20),
    "ng_fill": RGBColor(0xFD, 0xEC, 0xEC),
    "ng_text": RGBColor(0xB3, 0x26, 0x1E),
}


def _set_run(
    run,
    *,
    size: float,
    color: RGBColor,
    bold: bool = False,
) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: float,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
) -> object:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    frame.vertical_anchor = vertical
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = 0
    paragraph.space_after = 0
    paragraph.line_spacing = 1.0
    run = paragraph.add_run()
    run.text = text
    _set_run(run, size=size, color=color, bold=bold)
    return box


def _add_pill(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor,
    color: RGBColor,
    size: float = 8.0,
) -> None:
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    frame = pill.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    _set_run(run, size=size, color=color, bold=True)


def _add_card(
    slide,
    *,
    index: int,
    category: str,
    title: str,
    subtitle: str,
    left: float,
    top: float,
    width: float,
    height: float,
    accent: RGBColor,
    highlighted: bool = False,
    status_output: bool = False,
) -> None:
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["card_blue"] if highlighted else COLORS["card"]
    card.line.color.rgb = COLORS["card_edge"]
    card.line.width = Pt(1.0)

    accent_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left + 0.08),
        Inches(top + 0.02),
        Inches(width - 0.16),
        Inches(0.045),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    step = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(left + 0.13),
        Inches(top + 0.16),
        Inches(0.30),
        Inches(0.30),
    )
    step.fill.solid()
    step.fill.fore_color.rgb = accent
    step.line.fill.background()
    step_frame = step.text_frame
    step_frame.clear()
    step_frame.margin_left = 0
    step_frame.margin_right = 0
    step_frame.margin_top = 0
    step_frame.margin_bottom = 0
    step_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = step_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = f"{index:02d}"
    _set_run(run, size=7.5, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    _add_text(
        slide,
        category,
        left + 0.51,
        top + 0.16,
        width - 0.62,
        0.30,
        size=8.2,
        color=COLORS["muted"],
        bold=True,
    )
    _add_text(
        slide,
        title,
        left + 0.13,
        top + 0.51,
        width - 0.26,
        0.58,
        size=13.0,
        color=COLORS["text"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    if status_output:
        chip_width = 0.48
        chip_gap = 0.08
        chips_width = chip_width * 2 + chip_gap
        chip_left = left + (width - chips_width) / 2
        _add_pill(
            slide,
            "OK",
            chip_left,
            top + 1.19,
            chip_width,
            0.24,
            fill=COLORS["ok_fill"],
            color=COLORS["ok_text"],
            size=7.5,
        )
        _add_pill(
            slide,
            "NG",
            chip_left + chip_width + chip_gap,
            top + 1.19,
            chip_width,
            0.24,
            fill=COLORS["ng_fill"],
            color=COLORS["ng_text"],
            size=7.5,
        )
    else:
        _add_text(
            slide,
            subtitle,
            left + 0.10,
            top + 1.17,
            width - 0.20,
            0.26,
            size=7.8,
            color=COLORS["muted"],
            align=PP_ALIGN.CENTER,
        )


def _remove_other_slides(presentation: Presentation, keep_slide) -> None:
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        if presentation.part.related_part(slide_id.rId) is keep_slide.part:
            continue
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape.element)


def _validate(output: Path) -> None:
    if output.read_bytes()[:4] != b"PK\x03\x04":
        raise RuntimeError("生成文件没有有效的Office ZIP签名")
    with zipfile.ZipFile(output) as package:
        names = set(package.namelist())
        required = {
            "[Content_Types].xml",
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
        }
        missing = sorted(required - names)
        if missing:
            raise RuntimeError(f"PPTX缺少必要成员：{missing}")
        if not any(name.startswith("ppt/slides/slide") for name in names):
            raise RuntimeError("PPTX没有幻灯片XML")
    reopened = Presentation(output)
    if len(reopened.slides) != 1:
        raise RuntimeError(f"样本应为1页，实际为{len(reopened.slides)}页")
    text = "\n".join(
        shape.text for shape in reopened.slides[0].shapes if hasattr(shape, "text")
    )
    for expected in ("设备功能流程", "工位1", "分类下料"):
        if expected not in text:
            raise RuntimeError(f"样本缺少关键文字：{expected}")


def build(output: Path, *, overwrite: bool = False) -> Path:
    if output.exists() and not overwrite:
        raise FileExistsError(f"输出文件已存在：{output}")
    output.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(TEMPLATE)
    slide = presentation.slides[3]
    _remove_other_slides(presentation, slide)
    _clear_slide(slide)

    _add_text(
        slide,
        "2.2  设备功能流程",
        0.56,
        1.12,
        4.25,
        0.38,
        size=20.0,
        color=COLORS["text"],
        bold=True,
    )
    title_line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.47),
        Inches(1.60),
        Inches(4.08),
        Inches(0.012),
    )
    title_line.fill.solid()
    title_line.fill.fore_color.rgb = COLORS["connector"]
    title_line.line.fill.background()

    _add_pill(
        slide,
        "工业单主线 · 样本",
        10.82,
        1.16,
        1.91,
        0.32,
        fill=RGBColor(0xEA, 0xF1, 0xF7),
        color=COLORS["primary"],
        size=8.5,
    )

    _add_text(
        slide,
        "设备主流程",
        0.55,
        2.28,
        1.45,
        0.28,
        size=10.5,
        color=COLORS["primary"],
        bold=True,
    )
    _add_text(
        slide,
        "从产品输入到检测结果输出，保持单一阅读方向",
        2.00,
        2.28,
        4.50,
        0.28,
        size=8.5,
        color=COLORS["muted"],
    )

    nodes = [
        ("输入", "产品输入", "人工 / 自动上料", COLORS["primary"], True, False),
        ("输送", "输送定位", "到位并稳定", COLORS["transfer"], False, False),
        ("视觉工位", "工位1\n俯视检测", "相机采集与分析", COLORS["vision"], True, False),
        ("姿态切换", "产品翻转", "180°姿态切换", COLORS["transfer"], False, False),
        ("视觉工位", "工位2\n仰视检测", "第二视角复检", COLORS["vision"], True, False),
        ("结果处理", "结果判定", "汇总检测结果", COLORS["primary"], False, False),
        ("输出", "分类下料", "", COLORS["primary"], True, True),
    ]
    left = 0.52
    top = 2.84
    card_width = 1.45
    card_height = 1.58
    gap = 0.356
    for index, (category, title, subtitle, accent, highlighted, status) in enumerate(
        nodes, start=1
    ):
        card_left = left + (index - 1) * (card_width + gap)
        _add_card(
            slide,
            index=index,
            category=category,
            title=title,
            subtitle=subtitle,
            left=card_left,
            top=top,
            width=card_width,
            height=card_height,
            accent=accent,
            highlighted=highlighted,
            status_output=status,
        )
        if index < len(nodes):
            arrow_left = card_left + card_width + (gap - 0.21) / 2
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                Inches(arrow_left),
                Inches(top + 0.69),
                Inches(0.21),
                Inches(0.16),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["connector"]
            arrow.line.fill.background()

    note_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.52),
        Inches(4.87),
        Inches(12.25),
        Inches(0.67),
    )
    note_panel.fill.solid()
    note_panel.fill.fore_color.rgb = COLORS["background"]
    note_panel.line.color.rgb = COLORS["soft_line"]
    note_panel.line.width = Pt(0.8)
    note_items = (
        ("01", "统一方向", "所有节点从左到右阅读"),
        ("02", "信息分层", "工位名称优先，说明保持一行"),
        ("03", "状态克制", "OK / NG仅在输出位置显示"),
    )
    note_lefts = (0.82, 4.75, 8.68)
    for (number, title, detail), note_left in zip(note_items, note_lefts):
        _add_pill(
            slide,
            number,
            note_left,
            5.05,
            0.38,
            0.24,
            fill=RGBColor(0xDF, 0xE9, 0xF2),
            color=COLORS["primary"],
            size=7.2,
        )
        _add_text(
            slide,
            title,
            note_left + 0.50,
            4.97,
            0.95,
            0.22,
            size=8.7,
            color=COLORS["text"],
            bold=True,
        )
        _add_text(
            slide,
            detail,
            note_left + 0.50,
            5.19,
            2.65,
            0.22,
            size=7.6,
            color=COLORS["muted"],
        )

    _add_text(
        slide,
        "图 2.2  设备功能流程图",
        4.95,
        6.28,
        3.45,
        0.34,
        size=10.5,
        color=COLORS["text"],
        align=PP_ALIGN.CENTER,
    )

    presentation.save(output)
    _validate(output)
    return output.resolve()


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--overwrite", action="store_true")
    args = parser.parse_args()
    print(build(args.output.resolve(), overwrite=args.overwrite))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
