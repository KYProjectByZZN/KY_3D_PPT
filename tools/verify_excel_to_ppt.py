"""Apply the demo Excel rules and verify them against the NAT6704 PPT template."""

from __future__ import annotations

import argparse
import json
import sys
import tempfile
from pathlib import Path

from pptx import Presentation


PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from ppt_generator import (
    ExcelMappingRule,
    apply_excel_mappings,
    load_manifest,
    render_template,
)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "--excel", type=Path, default=PROJECT_ROOT / "examples" / "excel_mapping_demo.xlsx"
    )
    parser.add_argument(
        "--rules", type=Path, default=PROJECT_ROOT / "examples" / "excel_mapping_demo.json"
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=PROJECT_ROOT / "output" / "NAT6704_v2_Excel映射验证.pptx",
    )
    parser.add_argument("--overwrite", action="store_true")
    args = parser.parse_args()

    manifest_path = PROJECT_ROOT / "templates" / "NAT6704_v2.template.json"
    template_path = PROJECT_ROOT / "templates" / "冲压筒形壳体检测方案NAT6704_v2.pptx"
    base_data_path = PROJECT_ROOT / "examples" / "NAT6704_v2_test_data.json"
    manifest = load_manifest(manifest_path)
    raw_rules = json.loads(args.rules.read_text(encoding="utf-8"))
    rules = [ExcelMappingRule.from_dict(item) for item in raw_rules["mappings"]]
    mapped = apply_excel_mappings(
        args.excel,
        rules,
        slot_specs={slot["key"]: slot for slot in manifest.slots},
    )
    values = json.loads(base_data_path.read_text(encoding="utf-8"))
    values.update(mapped.values)
    with tempfile.TemporaryDirectory(prefix="ky_excel_ppt_") as temp_dir:
        data_path = Path(temp_dir) / "mapped.json"
        data_path.write_text(json.dumps(values, ensure_ascii=False), encoding="utf-8")
        render_template(
            template_path,
            manifest_path,
            data_path,
            args.output,
            overwrite=args.overwrite,
        )

    presentation = Presentation(args.output)
    title = next(shape.text for shape in presentation.slides[0].shapes if shape.shape_id == 7)
    parameters = next(
        shape.table for shape in presentation.slides[9].shapes if shape.shape_id == 3
    )
    inspections = next(
        shape.table for shape in presentation.slides[15].shapes if shape.shape_id == 12
    )
    if "Excel映射验证" not in title:
        raise RuntimeError("封面标题没有写入 Excel 映射值")
    if parameters.cell(0, 1).text != "36 pcs/min":
        raise RuntimeError("设备参数表没有写入 Excel 映射值")
    if inspections.cell(1, 1).text != "A面压痕":
        raise RuntimeError("检测项目表没有写入 Excel 映射值")
    print(
        f"verified={args.output.resolve()} slides={len(presentation.slides)} "
        f"mapped_slots={len(mapped.values)}"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
