"""Generate and structurally verify a real NAT6704 Excel-repeat PPTX."""

from __future__ import annotations

import json
import sys
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from ppt_generator import (
    ExcelModuleBinding,
    PptProject,
    ensure_project_modules,
    load_manifest,
    materialize_excel_modules,
    render_project,
    sha256_file,
)


TEMPLATE = PROJECT_ROOT / "templates" / "冲压筒形壳体检测方案NAT6704_v2.pptx"
MANIFEST = PROJECT_ROOT / "templates" / "NAT6704_v2.template.json"
DATA = PROJECT_ROOT / "examples" / "NAT6704_v2_test_data.json"
EXCEL = PROJECT_ROOT / "examples" / "module_repeat_demo.xlsx"
OUTPUT = PROJECT_ROOT / "output" / "NAT6704_v2_Excel模块重复验证.pptx"


def main() -> int:
    manifest = load_manifest(MANIFEST)
    original_hash = sha256_file(TEMPLATE)
    project = PptProject(
        project_name="Excel 模块重复验证",
        template_path=str(TEMPLATE),
        manifest_path=str(MANIFEST),
        output_path=str(OUTPUT),
        values=json.loads(DATA.read_text(encoding="utf-8")),
    )
    ensure_project_modules(project, manifest)
    source = next(
        item for item in project.modules if item.template_module_key == "equipment_overview"
    )
    binding = ExcelModuleBinding(
        source_module_id=source.id,
        source_path=str(EXCEL),
        sheet="设备模块",
        data_range="A2:B4",
        field_map={
            "设备标题": "equipment_title",
            "设备说明": "equipment_description",
        },
        module_name_field="设备标题",
    )
    generated = materialize_excel_modules(project, binding)
    render_project(project, OUTPUT, overwrite=True)

    presentation = Presentation(OUTPUT)
    expected_slides = 25  # 23 template pages - 1 source + 3 Excel copies
    if len(presentation.slides) != expected_slides:
        raise RuntimeError(
            f"Excel 模块重复页数应为 {expected_slides}，实际为 {len(presentation.slides)}"
        )
    with ZipFile(OUTPUT) as package:
        if package.testzip() is not None:
            raise RuntimeError("输出 PPTX ZIP 成员损坏")
        slide_xml_count = sum(
            name.startswith("ppt/slides/slide")
            and name.endswith(".xml")
            and "/_rels/" not in name
            for name in package.namelist()
        )
    if slide_xml_count != expected_slides:
        raise RuntimeError("输出 PPTX 含未引用或缺失的 Slide XML")
    if sha256_file(TEMPLATE) != original_hash:
        raise RuntimeError("原模板哈希发生变化")
    print(
        f"output={OUTPUT.resolve()} generated_modules={len(generated)} "
        f"slides={expected_slides} source_hash={original_hash}"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
